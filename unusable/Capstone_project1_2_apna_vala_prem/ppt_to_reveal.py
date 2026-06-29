"""
ppt_to_reveal.py — Convert a .pptx file to a reveal.js HTML presentation
                   with a WebSocket listener for PILOT voice navigation.

Usage:
    python ppt_to_reveal.py presentation.pptx
    python ppt_to_reveal.py presentation.pptx --output slides.html
"""

import sys
import argparse
import html as html_lib
from pathlib import Path


def extract_slides(pptx_path: str) -> list[dict]:
    """Extract text content from each slide."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)

    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        title = ""
        bullets = []

        for shape in slide.shapes:

    # Skip pictures, charts, tables, etc.
            if not getattr(shape, "has_text_frame", False):
                continue

            text = shape.text_frame.text.strip()

            if not text:
                continue

            is_title = False

            if shape.shape_type == 13:
                is_title = True

            elif shape.is_placeholder:
                try:
                    is_title = shape.placeholder_format.idx == 0
                except Exception:
                    pass

            if not title and is_title:
                title = text
            else:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        bullets.append(line)
        slides.append({
            "index": i,
            "title": title or f"Slide {i + 1}",
            "bullets": bullets,
        })

    return slides


def build_html(slides: list[dict], ws_port: int = 8001) -> str:
    """Build a full reveal.js HTML string with WebSocket nav listener."""

    # Build <section> blocks
    sections = []
    for slide in slides:
        title_safe = html_lib.escape(slide["title"])
        bullets_html = ""
        if slide["bullets"]:
            items = "".join(
                f"<li>{html_lib.escape(b)}</li>" for b in slide["bullets"]
            )
            bullets_html = f"<ul>{items}</ul>"

        sections.append(f"""
        <section>
            <h2>{title_safe}</h2>
            {bullets_html}
        </section>""")

    sections_html = "\n".join(sections)

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PILOT Slides</title>

    <!-- reveal.js CDN -->
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/reveal.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/theme/black.min.css">

    <style>
        /* PILOT voice indicator */
        #pilot-status {{
            position: fixed;
            bottom: 16px;
            right: 20px;
            background: rgba(0,0,0,0.6);
            color: #4ade80;
            font-family: monospace;
            font-size: 13px;
            padding: 6px 12px;
            border-radius: 6px;
            z-index: 9999;
            transition: opacity 0.3s;
        }}
        #pilot-status.flash {{
            background: #4ade80;
            color: #000;
        }}
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
            {sections_html}
        </div>
    </div>

    <!-- PILOT voice indicator -->
    <div id="pilot-status">🎙 PILOT connected</div>

    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/reveal.min.js"></script>
    <script>
        // ── Init reveal.js ────────────────────────────────────────
        Reveal.initialize({{
            hash: true,
            controls: true,
            progress: true,
            slideNumber: true,
            transition: 'slide',
        }});

        // ── WebSocket connection to PILOT slide_server ────────────
        const wsProtocol = window.location.protocol === "https:" ? "wss:" : "ws:";
        const WS_URL = `${{wsProtocol}}//${{window.location.host}}/ws`;
        let ws = null;
        const statusEl = document.getElementById("pilot-status");

        function flash(msg) {{
            statusEl.textContent = msg;
            statusEl.classList.add("flash");
            setTimeout(() => {{
                statusEl.classList.remove("flash");
                statusEl.textContent = "🎙 PILOT connected";
            }}, 1500);
        }}

        function connect() {{
            ws = new WebSocket(WS_URL);

            ws.onopen = () => {{
                console.log("[PILOT] WebSocket connected");
                statusEl.textContent = "🎙 PILOT connected";
            }};

            ws.onmessage = (event) => {{
                try {{
                    const cmd = JSON.parse(event.data);
                    console.log("[PILOT] Command received:", cmd);

                    switch (cmd.action) {{
                        case "next":
                            Reveal.next();
                            flash("▶ Next slide");
                            break;
                        case "prev":
                        case "previous":
                            Reveal.prev();
                            flash("◀ Previous slide");
                            break;
                        case "first":
                            Reveal.slide(0, 0);
                            flash("⏮ First slide");
                            break;
                        case "last":
                            const total = Reveal.getTotalSlides() - 1;
                            Reveal.slide(total, 0);
                            flash("⏭ Last slide");
                            break;
                        case "goto":
                            const idx = Math.max(0, (cmd.index || 1) - 1); // 1-based from voice
                            Reveal.slide(idx, 0);
                            flash(`↪ Slide ${{idx + 1}}`);
                            break;
                        default:
                            console.warn("[PILOT] Unknown action:", cmd.action);
                    }}
                }} catch(e) {{
                    console.error("[PILOT] Bad command:", e);
                }}
            }};

            ws.onclose = () => {{
                statusEl.textContent = "⚠️ PILOT disconnected — retrying...";
                setTimeout(connect, 2000); // auto-reconnect
            }};

            ws.onerror = (err) => {{
                console.error("[PILOT] WebSocket error:", err);
            }};
        }}

        connect();
    </script>
</body>
</html>"""


def main():
    parser = argparse.ArgumentParser(description="Convert .pptx to reveal.js HTML for PILOT")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument("--output", default="slides.html", help="Output HTML filename")
    parser.add_argument("--ws-port", type=int, default=8001, help="PILOT slide_server WebSocket port")
    args = parser.parse_args()

    print(f"[ppt_to_reveal] Reading {args.pptx}...")
    slides = extract_slides(args.pptx)
    print(f"[ppt_to_reveal] Extracted {len(slides)} slides")

    html = build_html(slides, ws_port=args.ws_port)

    output_path = Path(args.output)
    output_path.write_text(html, encoding="utf-8")
    print(f"[ppt_to_reveal] Written to {output_path.resolve()}")
    print(f"\nNext steps:")
    print(f"  1. python slide_server.py")
    print(f"  2. Open http://localhost:8001 in your browser")
    print(f"  3. Launch PILOT and say 'next slide', 'previous slide', 'go to slide 3'")


if __name__ == "__main__":
    main()