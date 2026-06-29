import os
import json
import urllib.request
import urllib.error
import asyncio
import traceback
from sqlalchemy import create_engine, text

# Retrieve config from environment
LLM_PROVIDER = os.getenv("PILOT_LLM_PROVIDER", "echo").lower()
LLM_MODEL = os.getenv("PILOT_LLM_MODEL", "llama3.2")
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://127.0.0.1:11434/api/chat")
OPENAI_BASE_URL = os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_MODEL = os.getenv("OPENAI_MODEL", LLM_MODEL)
TTS_SPEAK_URL = os.getenv("PILOT_TTS_SPEAK_URL", "http://127.0.0.1:8000/speak")

# Retrieve DATABASE_URL dynamically, fallback to local SQLite portable DB
DATABASE_URL = os.getenv("PILOT_DATABASE_URL", "sqlite:///pilot_voice.db")

BACKGROUND_SYSTEM_PROMPT = """You are the Background Processing Agent for PILOT.
You handle complex, multi-step tasks, deep calculations, code generation, and detailed analysis that cannot be run on the low-latency conversation front-end.
You must execute the requested task carefully, perform any logic, and respond with a highly detailed, helpful final voice report.
Keep your report friendly but precise. Since the report will be spoken via TTS, make it sound natural (e.g., avoid reading out complex syntax or raw markdown, but summarize nicely)."""

def post_json(url: str, payload: dict, headers: dict | None = None, timeout: float = 60.0) -> dict:
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        url,
        data=data,
        headers={
            "Content-Type": "application/json",
            **(headers or {}),
        },
        method="POST",
    )
    with urllib.request.urlopen(req, timeout=timeout) as response:
        body = response.read().decode("utf-8")
        return json.loads(body) if body else {}

# We can allow pilot_asr to register a custom send_to_tts callback to prevent echo loops
_send_to_tts_callback = None

def register_send_to_tts_callback(callback):
    global _send_to_tts_callback
    _send_to_tts_callback = callback

def send_to_tts(text: str):
    """Send voice output back to the user via the TTS server."""
    if _send_to_tts_callback:
        try:
            print(f"[BackgroundAgent] Routing to registered TTS callback: {text[:100]}...")
            _send_to_tts_callback(text)
            return
        except Exception as e:
            print(f"[BackgroundAgent] Registered TTS callback failed: {e}. Falling back to direct post.")

    try:
        print(f"[BackgroundAgent] Sending to TTS directly: {text[:100]}...")
        post_json(TTS_SPEAK_URL, {"text": text}, timeout=30.0)
    except Exception as e:
        print(f"[BackgroundAgent] TTS delivery failed: {e}")

class BackgroundAgent:
    def __init__(self):
        self.queue = asyncio.Queue()
        self.loop_task = None

    def start(self):
        """Starts the background worker loop."""
        self.loop_task = asyncio.create_task(self._worker_loop())
        print("[BackgroundAgent] Started background worker loop.")

    def stop(self):
        if self.loop_task:
            self.loop_task.cancel()
            print("[BackgroundAgent] Stopped background worker loop.")

    def submit_task(self, task_type: str, task_details: str):
        """Submit a task to be processed asynchronously."""
        self.queue.put_nowait({
            "type": task_type,
            "details": task_details
        })
        print(f"[BackgroundAgent] Task submitted: {task_type}")

    async def _worker_loop(self):
        while True:
            try:
                task = await self.queue.get()
                task_type = task["type"]
                details = task["details"]

                print(f"[BackgroundAgent] Executing task '{task_type}' with details: '{details}'")
                
                # Process the task asynchronously
                result = await self._execute_task_async(task_type, details)
                
                # Send the final result to TTS
                send_to_tts(result)
                
                self.queue.task_done()
            except asyncio.CancelledError:
                break
            except Exception as e:
                print(f"[BackgroundAgent] Error in worker loop: {e}")
                traceback.print_exc()

    async def _execute_task_async(self, task_type: str, details: str) -> str:
        """Asynchronously processes background tasks, routing to sync or async handlers."""
        if task_type == "MCP_TOOL_QUERY":
            return await self._handle_mcp_tool_query(details)
        elif task_type == "DATABASE_QUERY":
            return await asyncio.to_thread(self._handle_database_query, details)
        elif task_type == "WRITE_FILE":
            return await asyncio.to_thread(self._handle_write_file, details)
        elif task_type == "SYSTEM_CHECK":
            return await asyncio.to_thread(self._handle_system_check, details)
        elif task_type == "COMPLEX_CALCULATION":
            return await asyncio.to_thread(self._handle_complex_calculation, details)
        elif task_type == "SLIDE_CONTROL":
            return await self._handle_slide_control(details)
        else:
            return f"Received an unknown background task type: {task_type}. I was unable to process it."

    async def call_mcp_tool_async(self, server_cmd: str, server_args: list[str], tool_name: str, arguments: dict) -> dict:
        """
        Connects to an external MCP server running over stdio,
        initializes a session, calls the specified tool, and returns the result.
        """
        try:
            from mcp import ClientSession, StdioServerParameters
            from mcp.client.stdio import stdio_client
        except ImportError:
            # If the library is not installed, output instructions and return simulated data
            print("[BackgroundAgent] Note: 'mcp' Python SDK is not installed. To run real MCP servers, run: pip install mcp")
            return {
                "error": "mcp_not_installed",
                "message": "To run real MCP servers, please install the library via: pip install mcp"
            }

        server_params = StdioServerParameters(
            command=server_cmd,
            args=server_args,
            env=os.environ.copy()
        )

        try:
            print(f"[BackgroundAgent] Connecting to MCP Server: {server_cmd} {' '.join(server_args)}")
            async with stdio_client(server_params) as (read, write):
                async with ClientSession(read, write) as session:
                    print("[BackgroundAgent] Initializing MCP Session...")
                    await session.initialize()
                    
                    print(f"[BackgroundAgent] Invoking tool '{tool_name}' with args {arguments}...")
                    result = await session.call_tool(tool_name, arguments=arguments)
                    print(f"[BackgroundAgent] MCP invocation completed: {result}")
                    return {
                        "status": "success",
                        "content": [c.text for c in result.content if hasattr(c, 'text')]
                    }
        except Exception as e:
            print(f"[BackgroundAgent] MCP Error: {e}")
            return {"error": str(e)}

    # async def _handle_mcp_tool_query(self, details: str) -> str:
    #     print("[BackgroundAgent] Executing MCP Tool Query...")
        
    #     # Configure your real external MCP server details.
    #     # E.g. to connect to a flight or web search MCP server
    #     MCP_SERVER_COMMAND = os.getenv("PILOT_MCP_COMMAND", "node")
    #     # Split command line arguments safely
    #     raw_args = os.getenv("PILOT_MCP_ARGS", "")
    #     MCP_SERVER_ARGS = raw_args.split() if raw_args else []
    #     MCP_TOOL_NAME = os.getenv("PILOT_MCP_TOOL", "get_flights")
        
    #     # Parse origin and destination if mentioned
    #     origin = "Mumbai" if "mumbai" in details.lower() else "unknown"
    #     destination = "Delhi" if "delhi" in details.lower() else "unknown"
        
    #     mcp_arguments = {
    #         "origin": origin,
    #         "destination": destination,
    #         "query": details
    #     }
        
    #     mcp_response = {}
    #     if MCP_SERVER_ARGS:
    #         mcp_response = await self.call_mcp_tool_async(
    #             server_cmd=MCP_SERVER_COMMAND,
    #             server_args=MCP_SERVER_ARGS,
    #             tool_name=MCP_TOOL_NAME,
    #             arguments=mcp_arguments
    #         )
        
    #     if not mcp_response or "error" in mcp_response:
    #         print("[BackgroundAgent] Real-time MCP server is not configured or reachable. Simulating API response.")
    #         realtime_data = (
    #             "Found 3 direct flights from Mumbai to Delhi today. "
    #             "Air India flight 806 leaves at 2:15 PM for 5,200 Rupees. "
    #             "IndiGo flight 2043 leaves at 4:30 PM for 4,800 Rupees. "
    #             "Vistara flight 930 leaves at 6:00 PM for 6,100 Rupees."
    #         )
    #         report_prompt = (
    #             f"The user asked for real-time information: '{details}'. "
    #             f"The MCP server is not currently connected, but here is the simulated flight data: {realtime_data}. "
    #             f"Please synthesize this into a beautifully spoken final report. Also, add a small polite note "
    #             f"that this is currently simulated and can be hooked up to a real flight API MCP server anytime."
    #         )
    #     else:
    #         mcp_data = " ".join(mcp_response.get("content", []))
    #         report_prompt = (
    #             f"The user asked: '{details}'. We queried our external MCP server using '{MCP_TOOL_NAME}' "
    #             f"and received this raw tool output: {mcp_data}. "
    #             f"Please write a concise, friendly spoken report for the user."
    #         )
            
    #     return await asyncio.to_thread(self._call_background_llm, report_prompt)
    #     async def _handle_mcp_tool_query(self, details: str) -> str:
    #         print("[BackgroundAgent] Executing MCP Tool Query...")

    #         # New: attempt to extract service type, origin, destination
    #         services = {
    #             "flights": ["flight", "flights", "airport"],
    #             "hotels": ["hotel", "room", "lodging", "stay"],
    #             "trains": ["train", "rail", "coach", "irctc"],
    #             "cabs":   ["cab", "taxi", "ola", "uber", "rapido"]
    #         }
    #         service_type = "flights" # Default
    #         for svc, keywords in services.items():
    #             if any(kw in details.lower() for kw in keywords):
    #                 service_type = svc
    #                 break

    #         # Extract city names (basic heuristic)
    #         import re
    #         city_pattern = re.compile(r'\b([A-Z][a-z]+)\b')
    #         cities = city_pattern.findall(details)
    #         origin = cities[0] if cities else "unknown"
    #         destination = cities[1] if len(cities) > 1 else "unknown"

    #         # MCP Tool name (adapt to service_type)
    #         MCP_TOOL_NAME_MAP = {
    #             "flights": "get_flights",
    #             "hotels": "get_hotels",
    #             "trains": "get_trains",
    #             "cabs":   "get_cabs"
    #         }
    #         MCP_TOOL_NAME = MCP_TOOL_NAME_MAP.get(service_type, "get_flights")

    #         MCP_SERVER_COMMAND = os.getenv("PILOT_MCP_COMMAND", "node")
    #         raw_args = os.getenv("PILOT_MCP_ARGS", "")
    #         MCP_SERVER_ARGS = raw_args.split() if raw_args else []
            
    #         mcp_arguments = {
    #             "serviceType": service_type,
    #             "origin": origin,
    #             "destination": destination,
    #             "query": details
    #         }

    #         mcp_response = {}
    #         if MCP_SERVER_ARGS:
    #             mcp_response = await self.call_mcp_tool_async(
    #                 server_cmd=MCP_SERVER_COMMAND,
    #                 server_args=MCP_SERVER_ARGS,
    #                 tool_name=MCP_TOOL_NAME,
    #                 arguments=mcp_arguments
    #             )

    #         # Simulated fallback data per category
    #         fallback_data = {
    #             "flights": (
    #                 "Found 3 direct flights from {} to {} today. "
    #                 "Air India flight 806 leaves at 2:15 PM for 5,200 Rupees. "
    #                 "IndiGo flight 2043 leaves at 4:30 PM for 4,800 Rupees. "
    #                 "Vistara flight 930 leaves at 6:00 PM for 6,100 Rupees."
    #             ),
    #             "hotels": (
    #                 "Found 2 hotels available in {}. The Taj Hotel offers rooms from ₹8500. "
    #                 "Oberoi Hotel starts from ₹6200. Customer care for Taj is 1800-266-7646, and for Oberoi is 1800-102-2333."
    #             ),
    #             "trains": (
    #                 "There are 2 train options from {} to {} today. "
    #                 "Rajdhani Express departs at 11:30 AM, contact IRCTC helpline at 139. "
    #                 "Shatabdi departs at 6:45 PM, contact RailYatri at 8010500300."
    #             ),
    #             "cabs": (
    #                 "Available cab options from {} to {} include Ola with customer care 0120-3355335 and Uber with 080-4685-2190. "
    #                 "Estimated fares range from ₹900 to ₹1500 depending on vehicle type."
    #             ),
    #         }
    #         if not mcp_response or "error" in mcp_response:
    #             print("[BackgroundAgent] Real-time MCP server is not configured or reachable. Simulating API response.")
    #             key = service_type
    #             data_str = fallback_data[key].format(origin, destination)
    #             polite_note = "This is simulated info. Connect a real MCP server for live data."
    #             report_prompt = (
    #                 f"The user asked for real-time information: '{details}'. "
    #                 f"The MCP server is not currently connected, but here is simulated data: {data_str} "
    #                 f"{polite_note}. Please synthesize this into a friendly spoken report."
    #             )
    #         else:
    #             mcp_data = " ".join(mcp_response.get("content", []))
    #             report_prompt = (
    #                 f"The user asked: '{details}'. We queried our external MCP server using '{MCP_TOOL_NAME}' "
    #                 f"and received this output: {mcp_data}. Please write a concise, friendly spoken report for the user."
    #             )

    #         return await asyncio.to_thread(self._call_background_llm, report_prompt)
    async def _handle_mcp_tool_query(self, details: str) -> str:
        print("[BackgroundAgent] Executing MCP Tool Query...")

        
        import os
        import re
        import shlex
        import asyncio

        query = details.lower()

        # --------------------------------------------------
        # Detect Service Type
        # --------------------------------------------------

        services = {
            "flights": ["flight", "flights", "airport"],
            "hotels": ["hotel", "room", "stay", "lodging"],
            "trains": ["train", "rail", "irctc"],
            "cabs": ["cab", "taxi", "uber", "ola", "rapido"],
            "slides": [
                "next slide",
                "previous slide",
                "slide",
                "presentation",
                "go back",
                "next"
            ]
        }

        service_type = "flights"

        for svc, keywords in services.items():
            if any(keyword in query for keyword in keywords):
                service_type = svc
                break

        # --------------------------------------------------
        # Slide Navigation
        # --------------------------------------------------

        if service_type == "slides":

            if "next" in query:
                await self.slide_controller.send({
                    "action": "next"
                })
                return "Moving to the next slide."

            if "previous" in query or "back" in query:
                await self.slide_controller.send({
                    "action": "prev"
                })
                return "Moving to the previous slide."

            slide_match = re.search(r"slide\s+(\d+)", query)

            if slide_match:
                slide_number = int(slide_match.group(1))

                await self.slide_controller.send({
                    "action": "goto",
                    "slide": slide_number
                })

                return f"Moving to slide {slide_number}."

            return "Slide command received."

        # --------------------------------------------------
        # Extract Cities
        # --------------------------------------------------

        known_cities = {
            "mumbai",
            "delhi",
            "chennai",
            "bangalore",
            "hyderabad",
            "kolkata",
            "pune",
            "ahmedabad",
            "jaipur",
            "lucknow"
        }

        words = re.findall(r"[a-zA-Z]+", query)

        cities = [
            word.title()
            for word in words
            if word in known_cities
        ]

        origin = cities[0] if len(cities) > 0 else "Unknown"
        destination = cities[1] if len(cities) > 1 else "Unknown"

        # --------------------------------------------------
        # MCP Configuration
        # --------------------------------------------------

        MCP_SERVER_COMMAND = os.getenv(
            "PILOT_MCP_COMMAND",
            "node"
        )
        if MCP_SERVER_COMMAND == "nodeexport":
            MCP_SERVER_COMMAND = "node"

        raw_args = os.getenv(
            "PILOT_MCP_ARGS",
            ""
        )

        MCP_SERVER_ARGS = (
            shlex.split(raw_args)
            if raw_args
            else []
        )

        # Prioritize the PILOT_MCP_TOOL env var if set (e.g. search_flights_web)
        env_mcp_tool = os.getenv("PILOT_MCP_TOOL")

        if env_mcp_tool:
            # If PILOT_MCP_TOOL is configured, use it for all service types (flights, hotels, trains, cabs)
            # as customer-care-mcp (index2.js) exposes a single dynamic tool mapping
            MCP_TOOL_NAME = env_mcp_tool
        else:
            # Fall back to individual hardcoded tool names if env var is not set
            MCP_TOOL_MAP = {
                "flights": "get_flights",
                "hotels": "get_hotels",
                "trains": "get_trains",
                "cabs": "get_cabs"
            }
            MCP_TOOL_NAME = MCP_TOOL_MAP.get(
                service_type,
                "get_flights"
            )

        mcp_arguments = {
            "serviceType": service_type,
            "origin": origin,
            "destination": destination,
            "query": details
        }

        # --------------------------------------------------
        # Call MCP Server
        # --------------------------------------------------

        mcp_response = {}

        if MCP_SERVER_ARGS:
            mcp_response = await self.call_mcp_tool_async(
                server_cmd=MCP_SERVER_COMMAND,
                server_args=MCP_SERVER_ARGS,
                tool_name=MCP_TOOL_NAME,
                arguments=mcp_arguments
            )

        # --------------------------------------------------
        # Fallback Responses
        # --------------------------------------------------

        fallback_data = {
            "flights": (
                f"Found 3 flights from {origin} to {destination}. "
                f"Air India ₹5200, IndiGo ₹4800, Vistara ₹6100."
            ),

            "hotels": (
                f"Found hotels in {destination}. "
                f"Taj Hotel starts at ₹8500 and Oberoi at ₹6200."
            ),

            "trains": (
                f"Rajdhani Express and Shatabdi are available "
                f"between {origin} and {destination}."
            ),

            "cabs": (
                f"Ola and Uber are available from "
                f"{origin} to {destination}. "
                f"Estimated fare ranges from ₹900 to ₹1500."
            )
        }

        # --------------------------------------------------
        # Generate Report
        # --------------------------------------------------

        if not mcp_response or "error" in mcp_response:

            print(
                "[BackgroundAgent] MCP unavailable. "
                "Using fallback data."
            )

            report_prompt = (
                f"The user asked: '{details}'. "
                f"The MCP server is unavailable. "
                f"Use this simulated data: "
                f"{fallback_data[service_type]}. "
                f"Generate a short, friendly spoken response. "
                f"Mention that this is simulated data."
            )

        else:

            content = mcp_response.get(
                "content",
                []
            )

            if isinstance(content, list):
                mcp_data = " ".join(
                    str(x)
                    for x in content
                )
            else:
                mcp_data = str(content)

            report_prompt = (
                f"The user asked: '{details}'. "
                f"The MCP server returned: "
                f"{mcp_data}. "
                f"Generate a concise spoken report."
            )

        return await asyncio.to_thread(
            self._call_background_llm,
            report_prompt
        )

   

    def _handle_database_query(self, details: str) -> str:
        print("[BackgroundAgent] Executing Database Query...")
        
        # 1. Gather actual database stats
        db_info = ""
        try:
            engine = create_engine(DATABASE_URL)
            with engine.connect() as conn:
                # Check recordings count
                result = conn.execute(text("SELECT COUNT(*) FROM recordings"))
                count = result.scalar()
                
                # Check recent recordings
                recent_result = conn.execute(text("SELECT user_name, prompt_text, recording_number FROM recordings ORDER BY id DESC LIMIT 3"))
                recent = recent_result.fetchall()
                
                db_info = f"The pilot voice database is connected. There are currently {count} audio recordings saved. "
                if count > 0:
                    records_summary = ", ".join([f"{r[0]}'s recording number {r[2]}" for r in recent])
                    db_info += f"The most recent recordings are: {records_summary}."
        except Exception as e:
            print(f"[BackgroundAgent] Database connection failed: {e}. Falling back to filesystem stats.")
            # Fallback to checking the uploads/ directory
            try:
                uploads_dir = "uploads"
                if os.path.exists(uploads_dir):
                    files = os.listdir(uploads_dir)
                    webm_files = [f for f in files if f.endswith(".webm")]
                    db_info = (f"The database server was unreachable, but checking the uploads directory, "
                               f"I found {len(webm_files)} audio files saved in storage.")
                else:
                    db_info = "The database is unreachable and the uploads directory does not exist."
            except Exception as fs_err:
                db_info = f"Both database connection and filesystem fallback failed. The database error was {e}."

        # 2. Use LLM to compose a natural voice response combining the prompt and the gathered info
        report_prompt = f"The user asked: '{details}'. Here is the actual data gathered: {db_info}. Please write a concise natural voice report summarizing this."
        return self._call_background_llm(report_prompt)

    def _handle_write_file(self, details: str) -> str:
        print("[BackgroundAgent] Executing File Generation...")
        
        # 1. Ask LLM to generate the file content and specify the filename
        prompt = (f"The user wants to generate a file based on: '{details}'. "
                  f"Please write the file content. "
                  f"You must format your response as a JSON object with two fields:\n"
                  f"- 'filename': name of the file to create (relative path to current directory)\n"
                  f"- 'content': the actual string content of the file\n"
                  f"Provide ONLY the raw JSON object.")
        
        raw_json = self._call_background_llm(prompt, raw=True)
        filename = "generated_output.txt"
        content = ""
        
        try:
            # Extract JSON from LLM response
            import re
            json_match = re.search(r"\{.*\}", raw_json, re.DOTALL)
            if json_match:
                parsed = json.loads(json_match.group(0))
                filename = parsed.get("filename", filename)
                content = parsed.get("content", "")
            else:
                content = raw_json
        except Exception as e:
            print(f"[BackgroundAgent] Failed to parse generated file JSON: {e}")
            content = raw_json

        # 2. Write the file safely to the workspace
        try:
            # Prevent directory traversal attacks
            safe_filename = os.path.basename(filename)
            # Support generating inside BackgroundLLM/ or FrontLLM/ if requested
            if "background" in filename.lower():
                safe_path = os.path.join("BackgroundLLM", safe_filename)
            elif "front" in filename.lower():
                safe_path = os.path.join("FrontLLM", safe_filename)
            else:
                safe_path = safe_filename

            with open(safe_path, "w", encoding="utf-8") as f:
                f.write(content)
            
            return f"I have successfully created and written the file {safe_path} in the workspace. The file contains {len(content)} characters of text."
        except Exception as e:
            return f"I generated the file content but encountered an error writing it to disk. The error was {e}."

    def _handle_system_check(self, details: str) -> str:
        print("[BackgroundAgent] Executing System Check...")
        
        # Gather directory listing
        try:
            files = os.listdir(".")
            # Filter and summarize list
            python_files = [f for f in files if f.endswith(".py")]
            folders = [f for f in files if os.path.isdir(f) and not f.startswith(".")]
            
            summary = f"In the current directory, I found {len(python_files)} python files, including {', '.join(python_files[:5])}. "
            if folders:
                summary += f"The main folders are {', '.join(folders)}."
        except Exception as e:
            summary = f"Failed to perform filesystem check: {e}"

        prompt = f"The user asked: '{details}'. Here is the directory listing/system stats summary: {summary}. Please write a concise natural voice report explaining this."
        return self._call_background_llm(prompt)

    def _get_latest_pptx(self) -> str | None:
        """Finds the path to the latest uploaded .pptx file in the presentations folder."""
        folder = "presentations"
        if not os.path.exists(folder):
            return None
        files = [os.path.join(folder, f) for f in os.listdir(folder) if f.endswith(".pptx")]
        if not files:
            return None
        return max(files, key=os.path.getmtime)

    def _get_slide_content(self, index: int) -> str | None:
        """Extracts text content from a specific slide of the latest uploaded PPTX directly."""
        pptx_path = self._get_latest_pptx()
        if not pptx_path:
            return None
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            if 1 <= index <= len(prs.slides):
                slide = prs.slides[index - 1]
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = paragraph.text.strip()
                            if line:
                                text_runs.append(line)
                return " ".join(text_runs)
        except Exception as e:
            print(f"[BackgroundAgent] Error reading PPTX slide {index}: {e}")
        return None

    def _get_all_pptx_content(self) -> str | None:
        """Extracts the entire textual presentation structure from the latest uploaded PPTX directly."""
        pptx_path = self._get_latest_pptx()
        if not pptx_path:
            return None
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            all_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = paragraph.text.strip()
                            if line:
                                slide_text.append(line)
                if slide_text:
                    all_text.append(f"Slide {i + 1}: " + " ".join(slide_text))
            return "\n".join(all_text)
        except Exception as e:
            print(f"[BackgroundAgent] Error reading entire PPTX structure: {e}")
        return None

    def _get_slide_image_paths(self, index: int) -> list[str]:
        """Finds any extracted image paths matching the current slide index."""
        folder = os.path.join("presentations", "images")
        if not os.path.exists(folder):
            return []
        matching_files = []
        try:
            for f in os.listdir(folder):
                if f"_slide_{index}_" in f:
                    matching_files.append(os.path.join(folder, f))
        except Exception as e:
            print(f"[BackgroundAgent] Error scanning presentation images: {e}")
        return matching_files

    def _handle_complex_calculation(self, details: str) -> str:
        print("[BackgroundAgent] Executing Complex Calculation...")
        
        # Let's perform some local Python execution or smart analysis based on instructions
        prompt = f"Please perform this complex task, reasoning, or math query: '{details}'. Compute the answer and present it in a clear, natural voice report."
        return self._call_background_llm(prompt)

    async def _handle_slide_control(self, details: str) -> str:
        """
        Parses a slide navigation command from natural language
        and sends it to slide_server via HTTP POST.
 
        Supported intents:
          "next slide"            → {"action": "next"}
          "previous / back"       → {"action": "prev"}
          "first slide"           → {"action": "first"}
          "last slide"            → {"action": "last"}
          "go to slide 4"         → {"action": "goto", "index": 4}
          "slide 2" / "slide two" → {"action": "goto", "index": 2}
        """
        import re
 
        SLIDE_SERVER_URL = os.getenv("PILOT_SLIDE_SERVER", "http://127.0.0.1:8001")
        text = details.lower().strip()
 
        word_to_num = {
            "one": 1, "two": 2, "three": 3, "four": 4, "five": 5,
            "six": 6, "seven": 7, "eight": 8, "nine": 9, "ten": 10,
            "eleven": 11, "twelve": 12, "thirteen": 13, "fourteen": 14,
            "fifteen": 15, "sixteen": 16, "seventeen": 17, "eighteen": 18,
            "nineteen": 19, "twenty": 20,
        }

        # ── Check if the user is asking a question rather than navigating ──
        is_question = any(q in text for q in ("what", "why", "how", "explain", "describe", "tell", "read", "show", "summarize", "about", "content"))
        if is_question:
            # Extract slide number/word if mentioned
            num_match = re.search(r"slide\s+(\d+)", text)
            word_match = None
            for word, num in word_to_num.items():
                if re.search(rf"\bslide\s+{word}\b", text):
                    word_match = num
                    break

            slide_index = None
            if num_match:
                slide_index = int(num_match.group(1))
            elif word_match:
                slide_index = word_match

            # If a specific slide is requested, answer from slide text + slide images or fallback
            if slide_index is not None:
                content = self._get_slide_content(slide_index)
                image_paths = self._get_slide_image_paths(slide_index)
                
                if image_paths:
                    print(f"[BackgroundAgent] Slide {slide_index} has {len(image_paths)} embedded figures. Analyzing using Vision LLM capability.")
                    slide_text_part = content if content else ""
                    prompt = (
                        f"The user is asking: '{details}'. I have found {len(image_paths)} figures/diagrams embedded directly "
                        f"in slide {slide_index} of the presentation, along with this text: '{slide_text_part}'.\n"
                        f"Please analyze both the slide text and the attached figures/diagrams carefully, and write a friendly, "
                        f"clear, and comprehensive spoken response answering the user's question directly."
                    )
                    return self._call_background_llm_with_images(prompt, image_paths)
                elif content:
                    print(f"[BackgroundAgent] Extracted content from slide {slide_index}, answering using slide text.")
                    prompt = (
                        f"The user is asking: '{details}'. Here is the text content from slide {slide_index} "
                        f"of the presentation:\n'{content}'\n"
                        f"Please provide a friendly, clear, and concise spoken response answering the user's question directly based on this slide content."
                    )
                    return self._call_background_llm(prompt)
                else:
                    print(f"[BackgroundAgent] Slide {slide_index} not loaded or out of bounds. Answering from general knowledge.")
                    prompt = (
                        f"The user is asking: '{details}'. There are no presentation slides currently loaded or the slide "
                        f"is empty, so please answer their question directly using your general knowledge of the presentation topic "
                        f"(Biosorption & Bioaccumulation, heavy metals, bioremediation, etc.) in a friendly, concise spoken response."
                    )
                    return self._call_background_llm(prompt)
            else:
                # General presentation question
                all_ppt_content = self._get_all_pptx_content()
                if all_ppt_content:
                    print(f"[BackgroundAgent] General presentation question. Summarizing from direct PPTX content.")
                    prompt = (
                        f"The user is asking a general question about the presentation: '{details}'. "
                        f"Here is the full textual layout extracted from the PowerPoint presentation slides:\n"
                        f"\"\"\"\n{all_ppt_content}\n\"\"\"\n"
                        f"Please provide a friendly, clear, and concise spoken response answering their question directly "
                        f"based on the extracted presentation content."
                    )
                    return self._call_background_llm(prompt)
                else:
                    print(f"[BackgroundAgent] General topic question. Answering from general knowledge.")
                    prompt = (
                        f"The user is asking: '{details}'. Please provide a friendly and concise spoken response answering "
                        f"their question using your general knowledge of the presentation topic (Biosorption & Bioaccumulation, "
                        f"heavy metals, bioremediation, etc.). Keep it clear, friendly, and natural."
                    )
                    return self._call_background_llm(prompt)

        # ── Parse navigation intent ───────────────────────────
        command = {"action": "next"}   # safe default
 
        if any(w in text for w in ("previous", "back", "last time", "go back")):
            command = {"action": "prev"}
 
        elif any(w in text for w in ("first", "beginning", "start", "restart")):
            command = {"action": "first"}
 
        elif any(w in text for w in ("last", "final", "end")):
            command = {"action": "last"}
 
        else:
            # "go to slide 4" / "slide four" / "jump to 3"
            num_match = re.search(r"slide\s+(\d+)", text)
            word_match = None
            for word, num in word_to_num.items():
                if re.search(rf"\bslide\s+{word}\b", text):
                    word_match = num
                    break
 
            if num_match:
                command = {"action": "goto", "index": int(num_match.group(1))}
            elif word_match:
                command = {"action": "goto", "index": word_match}
            else:
                # Default: forward
                command = {"action": "next"}
 
        # ── Send to slide_server ──────────────────────────────
        print(f"[BackgroundAgent] Slide command: {command}")
        try:
            result = post_json(
                f"{SLIDE_SERVER_URL}/slide/command",
                command,
                timeout=5.0
            )
            clients = result.get("clients_notified", 0)
            action = command["action"]
 
            if clients == 0:
                print("[BackgroundAgent] Note: Slide command sent but no browser is connected to slide server.")
 
            # Natural TTS response per action
            responses = {
                "next":  "Moving to the next slide.",
                "prev":  "Going back to the previous slide.",
                "first": "Jumping to the first slide.",
                "last":  "Jumping to the last slide.",
                "goto":  f"Going to slide {command.get('index', '?')}.",
            }
            return responses.get(action, "Slide navigated.")
 
        except Exception as e:
            print(f"[BackgroundAgent] Slide command failed: {e}")
            return "I couldn't reach the slide server. Make sure slide_server.py is running."
 

    def _call_background_llm(self, prompt: str, raw: bool = False) -> str:
        """Call the background LLM (using higher capability, longer timeout)."""
        if LLM_PROVIDER in {"echo", "mock", "none"}:
            return f"Background agent completed processing: {prompt}"

        system_prompt = BACKGROUND_SYSTEM_PROMPT
        if raw:
            system_prompt += " Return ONLY the requested JSON object without any formatting/markdown blocks."

        if LLM_PROVIDER == "ollama":
            payload = {
                "model": LLM_MODEL,
                "stream": False,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt},
                ],
                "options": {
                    "temperature": 0.3
                }
            }
            try:
                result = post_json(OLLAMA_URL, payload, timeout=180.0)
                return result.get("message", {}).get("content", "").strip()
            except Exception as e:
                return f"Background LLM call failed: {e}"

        elif LLM_PROVIDER == "openai":
            if not OPENAI_API_KEY:
                return "Background LLM call failed: OpenAI API key not configured."
            payload = {
                "model": OPENAI_MODEL,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt},
                ],
                "temperature": 0.3,
            }
            try:
                result = post_json(
                    f"{OPENAI_BASE_URL.rstrip('/')}/chat/completions",
                    payload,
                    headers={"Authorization": f"Bearer {OPENAI_API_KEY}"},
                    timeout=180.0,
                )
                return result["choices"][0]["message"]["content"].strip()
            except Exception as e:
                return f"Background LLM call failed: {e}"

        return f"Unknown LLM provider: {LLM_PROVIDER}"

    def _call_background_llm_with_images(self, prompt: str, image_paths: list[str], raw: bool = False) -> str:
        """Call the background LLM with vision-enabled payload to analyze figures directly."""
        if LLM_PROVIDER in {"echo", "mock", "none"}:
            return f"Background agent completed processing with images: {prompt}"

        system_prompt = BACKGROUND_SYSTEM_PROMPT
        if raw:
            system_prompt += " Return ONLY the requested JSON object without any formatting/markdown blocks."

        import base64
        base64_images = []
        for path in image_paths:
            try:
                if os.path.exists(path):
                    with open(path, "rb") as img_file:
                        b64 = base64.b64encode(img_file.read()).decode("utf-8")
                        base64_images.append((path, b64))
            except Exception as e:
                print(f"[BackgroundAgent] Error encoding image {path}: {e}")

        if not base64_images:
            # Fallback to standard call if encoding failed or no files found
            return self._call_background_llm(prompt, raw=raw)

        if LLM_PROVIDER == "openai":
            if not OPENAI_API_KEY:
                return "Background LLM call failed: OpenAI API key not configured."
            
            content_list = [{"type": "text", "text": prompt}]
            for path, b64 in base64_images:
                ext = path.split(".")[-1].lower()
                mime = f"image/{ext}" if ext in {"png", "jpg", "jpeg", "gif", "webp"} else "image/png"
                content_list.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{mime};base64,{b64}"
                    }
                })

            payload = {
                "model": "gpt-4o-mini" if "mini" in OPENAI_MODEL else OPENAI_MODEL,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": content_list},
                ],
                "temperature": 0.3,
            }
            try:
                result = post_json(
                    f"{OPENAI_BASE_URL.rstrip('/')}/chat/completions",
                    payload,
                    headers={"Authorization": f"Bearer {OPENAI_API_KEY}"},
                    timeout=180.0,
                )
                return result["choices"][0]["message"]["content"].strip()
            except Exception as e:
                return f"Background LLM Vision call failed: {e}"

        elif LLM_PROVIDER == "ollama":
            images_list = [b64 for _, b64 in base64_images]
            payload = {
                "model": LLM_MODEL,
                "stream": False,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {
                        "role": "user",
                        "content": prompt,
                        "images": images_list
                    },
                ],
                "options": {
                    "temperature": 0.3
                }
            }
            try:
                result = post_json(OLLAMA_URL, payload, timeout=180.0)
                return result.get("message", {}).get("content", "").strip()
            except Exception as e:
                return f"Background LLM Ollama Vision call failed: {e}"

        return f"Unknown LLM provider: {LLM_PROVIDER}"

# Singleton instance
background_agent = BackgroundAgent()
