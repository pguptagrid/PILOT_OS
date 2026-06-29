


"""
PILOT — TTS Streaming Server  (with interrupt support)
Stack: FastAPI WebSocket + Kokoro ONNX → chunked PCM → browser playback

Interrupt protocol:
  Client → server : "__INTERRUPT__"  (text frame, any time during playback)
  Server           : stops Kokoro stream, sends "__INTERRUPTED__" back
  Client           : flushes audio queue, stops playback immediately

All other protocol unchanged:
  Client → server : plain text to speak
  Server → client : binary PCM chunks (Int16, 24 kHz, mono)
  Server → client : "__END__" when utterance finishes naturally
"""

import asyncio
import json
import os

import numpy as np
from fastapi import FastAPI, WebSocket, WebSocketDisconnect
from fastapi.responses import HTMLResponse
from kokoro_onnx import Kokoro
from pydantic import BaseModel


# ─────────────────────────────────────────────────────────────
#  CONFIG
# ─────────────────────────────────────────────────────────────

MODEL_PATH  = os.getenv("KOKORO_MODEL",  "kokoro-v1.0.onnx")
VOICES_PATH = os.getenv("KOKORO_VOICES", "voices-v1.0.bin")
VOICE       = os.getenv("KOKORO_VOICE",  "af_heart")
SPEED       = float(os.getenv("KOKORO_SPEED", "1.0"))
SAMPLE_RATE = 24000
CHUNK_BYTES = 4096


# ─────────────────────────────────────────────────────────────
#  MODEL
# ─────────────────────────────────────────────────────────────

kokoro = Kokoro(MODEL_PATH, VOICES_PATH)
print(f"[TTS] Kokoro loaded — voice={VOICE}  rate={SAMPLE_RATE}Hz")


# ─────────────────────────────────────────────────────────────
#  PCM HELPERS
# ─────────────────────────────────────────────────────────────

def float32_to_int16_bytes(audio: np.ndarray) -> bytes:
    peak = float(np.max(np.abs(audio)))
    if peak > 1e-6:
        audio = audio * (0.95 / peak)
    return np.clip(audio, -1.0, 1.0).astype(np.float32).__mul__(32767).astype(np.int16).tobytes()


def iter_chunks(pcm_bytes: bytes, chunk_size: int = CHUNK_BYTES):
    for i in range(0, len(pcm_bytes), chunk_size):
        yield pcm_bytes[i : i + chunk_size]


# ─────────────────────────────────────────────────────────────
#  APP + SHARED STATE
# ─────────────────────────────────────────────────────────────

app = FastAPI(title="PILOT TTS Server")


class SpeakRequest(BaseModel):
    text: str


playback_clients: list[WebSocket] = []
playback_lock = asyncio.Lock()

transcript_clients: set[WebSocket] = set()
transcript_lock = asyncio.Lock()


# ─────────────────────────────────────────────────────────────
#  TRANSCRIPT BROADCAST
# ─────────────────────────────────────────────────────────────

async def broadcast_transcript(text: str, source: str = "llm") -> None:
    if not transcript_clients:
        return
    payload = json.dumps({"type": "transcript", "source": source, "text": text})
    dead: list[WebSocket] = []
    async with transcript_lock:
        for ws in list(transcript_clients):
            try:
                await ws.send_text(payload)
            except Exception:
                dead.append(ws)
        for ws in dead:
            transcript_clients.discard(ws)


# ─────────────────────────────────────────────────────────────
#  INTERRUPTIBLE TTS STREAM
#
#  interrupt_event: asyncio.Event — set it from anywhere to
#  stop the Kokoro stream mid-sentence for this connection.
# ─────────────────────────────────────────────────────────────

async def send_tts_interruptible(
    ws: WebSocket,
    text: str,
    interrupt_event: asyncio.Event,
) -> tuple[int, int, bool]:
    """
    Stream TTS to ws, checking interrupt_event between every chunk.

    Returns (chunk_count, bytes_sent, was_interrupted).
    If interrupted:
      - stops the Kokoro generator immediately
      - sends "__INTERRUPTED__" so the client can flush its queue
    """
    chunk_count   = 0
    bytes_sent    = 0
    interrupted   = False

    async for audio_chunk, _sr in kokoro.create_stream(
        text, voice=VOICE, speed=SPEED, lang="en-us"
    ):
        # Check interrupt flag before sending each Kokoro chunk
        if interrupt_event.is_set():
            interrupted = True
            break

        pcm_bytes = float32_to_int16_bytes(audio_chunk)

        for piece in iter_chunks(pcm_bytes):
            # Check again between every 4 KB piece
            if interrupt_event.is_set():
                interrupted = True
                break
            await ws.send_bytes(piece)
            chunk_count += 1
            bytes_sent  += len(piece)
            await asyncio.sleep(0)   # yield to event loop

        if interrupted:
            break

    if interrupted:
        await ws.send_text("__INTERRUPTED__")
        print(f"[TTS] ⚡ Interrupted after {chunk_count} chunks")
    else:
        await ws.send_text("__END__")

    interrupt_event.clear()   # reset for the next utterance
    return chunk_count, bytes_sent, interrupted


# ─────────────────────────────────────────────────────────────
#  BROADCAST TTS  (for /speak REST endpoint)
# ─────────────────────────────────────────────────────────────

# Global interrupt event for broadcast — one for all passive listeners
_broadcast_interrupt = asyncio.Event()


async def broadcast_tts(text: str) -> dict:
    _broadcast_interrupt.clear()
    async with playback_lock:
        if not playback_clients:
            return {"status": "no_clients", "clients": 0}

        # Select only the most recently connected client to avoid double voice/echo from multiple tabs
        active_ws = playback_clients[-1]
        dead, total_chunks, total_bytes = [], 0, 0

        try:
            chunks, sent, _ = await send_tts_interruptible(
                active_ws, text, _broadcast_interrupt
            )
            total_chunks += chunks
            total_bytes  += sent
        except Exception:
            dead.append(active_ws)

        for ws in dead:
            if ws in playback_clients:
                playback_clients.remove(ws)

        return {
            "status": "ok",
            "clients": len(playback_clients),
            "chunks": total_chunks,
            "bytes": total_bytes,
        }


# ─────────────────────────────────────────────────────────────
#  ENDPOINTS
# ─────────────────────────────────────────────────────────────

@app.get("/health")
async def health():
    return {
        "status": "ok",
        "voice": VOICE,
        "sample_rate": SAMPLE_RATE,
        "playback_clients": len(playback_clients),
        "transcript_clients": len(transcript_clients),
    }


from fastapi.staticfiles import StaticFiles

@app.get("/", response_class=HTMLResponse)
async def index():
    html_path = os.path.join(os.path.dirname(__file__), "tts_client.html")
    with open(html_path) as f:
        return f.read()

# Serve PPT files statically from the presentations folder
app.mount("/presentations", StaticFiles(directory="presentations"), name="presentations")


@app.post("/speak")
async def speak(req: SpeakRequest):
    text = req.text.strip()
    if not text:
        return {"status": "ignored", "reason": "empty_text"}
    print(f"[TTS] /speak: '{text[:80]}'")
    await broadcast_transcript(text, source="llm")
    return await broadcast_tts(text)


@app.post("/interrupt")
async def interrupt():
    """REST endpoint — set the global interrupt flag for broadcast listeners."""
    _broadcast_interrupt.set()
    print("[TTS] /interrupt called")
    return {"status": "interrupted"}


@app.websocket("/listen")
async def listen_socket(ws: WebSocket):
    """
    Passive playback + interrupt.
    Client can send "__INTERRUPT__" at any time to stop current speech.
    """
    await ws.accept()
    playback_clients.append(ws)
    print(f"[TTS] Playback client connected: {ws.client}")

    try:
        while True:
            msg = await ws.receive_text()
            if msg == "__INTERRUPT__":
                _broadcast_interrupt.set()
                print(f"[TTS] ⚡ Interrupt from passive client: {ws.client}")
    except WebSocketDisconnect:
        print(f"[TTS] Playback client disconnected: {ws.client}")
    finally:
        if ws in playback_clients:
            playback_clients.remove(ws)


@app.websocket("/transcript")
async def transcript_socket(ws: WebSocket):
    """Clients receive JSON text events when LLM speaks."""
    await ws.accept()
    transcript_clients.add(ws)
    print(f"[TTS] Transcript client: {ws.client}")
    try:
        while True:
            await ws.receive_text()
    except WebSocketDisconnect:
        pass
    finally:
        transcript_clients.discard(ws)


@app.websocket("/tts")
async def tts_socket(ws: WebSocket):
    """
    Interactive TTS WebSocket — one utterance at a time.

    Client → server:
      - any text         → synthesise and stream
      - "__INTERRUPT__"  → stop current stream immediately

    Server → client:
      - binary frames    → PCM chunks (Int16, 24 kHz)
      - "__END__"        → utterance finished naturally
      - "__INTERRUPTED__"→ utterance was cut short
    """
    await ws.accept()
    print(f"[TTS] /tts client: {ws.client}")

    interrupt_event = asyncio.Event()
    speak_task: asyncio.Task | None = None

    async def _receiver():
        """Reads incoming frames; sets interrupt on "__INTERRUPT__"."""
        nonlocal speak_task
        while True:
            msg = await ws.receive_text()
            msg = msg.strip()

            if msg == "__INTERRUPT__":
                interrupt_event.set()
                print(f"[TTS] ⚡ Interrupt from /tts client: {ws.client}")

            elif msg:
                # New utterance — interrupt any current stream first
                interrupt_event.set()
                if speak_task and not speak_task.done():
                    await asyncio.sleep(0.05)   # let current chunk finish

                interrupt_event.clear()
                print(f"[TTS] Synthesising: '{msg[:80]}'")
                await broadcast_transcript(msg, source="llm")

                speak_task = asyncio.create_task(
                    send_tts_interruptible(ws, msg, interrupt_event)
                )
                chunks, sent, interrupted = await speak_task
                print(
                    f"[TTS] {'Interrupted' if interrupted else 'Done'} — "
                    f"{chunks} chunks, {sent / 1024:.1f} KB"
                )

    try:
        await _receiver()
    except WebSocketDisconnect:
        print(f"[TTS] /tts client disconnected: {ws.client}")
    except Exception as e:
        print(f"[TTS] Error: {e}")
        await ws.close(code=1011)



# from flask import request, jsonify
# from fastapi import UploadFile, File
# import os

# UPLOAD_FOLDER = "presentations"
# os.makedirs(UPLOAD_FOLDER, exist_ok=True)
# @app.post("/upload_ppt")
# async def upload_ppt(file: UploadFile = File(...)):
#     path = os.path.join(
#         UPLOAD_FOLDER,
#         file.filename
#     )

#     contents = await file.read()

#     with open(path, "wb") as f:
#         f.write(contents)

#     print(f"[PPT] Uploaded: {file.filename}")

#     return {
#         "status": "success",
#         "filename": file.filename,
#         "path": path
#     }

from fastapi import UploadFile, File
import subprocess
import asyncio

UPLOAD_FOLDER = "presentations"
SLIDES_OUTPUT = "slides.html"
SLIDE_SERVER_PORT = int(os.getenv("PILOT_SLIDE_PORT", "8001"))
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Global reference so we don't spawn multiple slide servers
_slide_server_proc: subprocess.Popen | None = None


@app.post("/upload_ppt")
async def upload_ppt(ppt: UploadFile = File(...)):   # ← "ppt" matches HTML field name
    # 1. Save the uploaded file
    save_path = os.path.join(UPLOAD_FOLDER, ppt.filename)
    contents = await ppt.read()
    with open(save_path, "wb") as f:
        f.write(contents)
    print(f"[PPT] Saved PPTX directly: {save_path}")

    # Ensure the static presentations/images directory exists
    images_folder = os.path.join(UPLOAD_FOLDER, "images")
    os.makedirs(images_folder, exist_ok=True)

    # Extract all slides text, metadata, and embedded figures / pictures
    slides_data = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import uuid
        
        file_uuid = uuid.uuid4().hex[:8]
        prs = Presentation(save_path)
        
        for i, slide in enumerate(prs.slides):
            title = f"Slide {i + 1}"
            bullets = []
            slide_images = []
            
            for shape in slide.shapes:
                # 1. Extract slide texts
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # Try to identify titles vs bullets
                        is_title = False
                        if shape.is_placeholder:
                            try:
                                is_title = shape.placeholder_format.idx == 0
                            except Exception:
                                pass
                        
                        if is_title:
                            title = text
                        else:
                            for paragraph in shape.text_frame.paragraphs:
                                line = paragraph.text.strip()
                                if line:
                                    bullets.append(line)
                                    
                # 2. Extract slide embedded figures / pictures directly from XML shape layers
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        ext = image.ext or "png"
                        image_filename = f"img_{file_uuid}_slide_{i+1}_{shape.shape_id}.{ext}"
                        image_save_path = os.path.join(images_folder, image_filename)
                        with open(image_save_path, "wb") as img_f:
                            img_f.write(image_bytes)
                        
                        # Store local static serving URL for front-end rendering
                        slide_images.append(f"/presentations/images/{image_filename}")
                        print(f"[PPT] Extracted embedded figure slide {i+1}: {image_filename}")
                    except Exception as img_err:
                        print(f"[PPT] Error extracting embedded picture shape on slide {i+1}: {img_err}")
            
            slides_data.append({
                "index": i + 1,
                "title": title,
                "bullets": bullets,
                "images": slide_images
            })
    except Exception as e:
        print(f"[PPT] Direct slides metadata extraction failed: {e}")

    # Return slide_server url and parsed slide cards for rich UI visualization
    return {
        "status": "success",
        "filename": ppt.filename,
        "slides_url": f"http://localhost:{SLIDE_SERVER_PORT}",
        "slide_count": len(slides_data) if slides_data else _count_slides(save_path),
        "slides": slides_data
    }


def _count_slides(pptx_path: str) -> int:
    try:
        from pptx import Presentation
        return len(Presentation(pptx_path).slides)
    except Exception:
        return -1
# .venv/bin/python -m uvicorn tts_server:app --reload
# export PILOT_LLM_PROVIDER=ollama
# export PILOT_LLM_MODEL=mistral:latest
# export OLLAMA_URL=http://127.0.0.1:11434/api/chat
# export PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak
# PILOT_LLM_PROVIDER=ollama PILOT_LLM_MODEL=mistral:latest OLLAMA_URL=http://127.0.0.1:11434/api/chat PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak .venv/bin/python pilot_asr.py
# .venv/bin/python pilot_asr.py
# PILOT_LLM_PROVIDER=ollama PILOT_LLM_MODEL=mistral:latest OLLAMA_URL=http://127.0.0.1:11434/api/chat PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak .venv/bin/python pilot_asr.py


# ollama serve                          .venv/bin/python -m uvicorn tts_server:app --reload                  
# ollama pull mistral:latest


# export PILOT_LLM_PROVIDER=ollamaexport PILOT_LLM_MODEL=mistral:latest                                                   
# export OLLAMA_URL=http://127.0.0.1:11434/api/chat
# export PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak


#  export KOKORO_MODEL=kokoro-v1.0.onnx                                   
# export KOKORO_VOICES=voices-v1.0.bin

# 
# export OLLAMA_URL=http://127.0.0.1:11434/api/chat 

#  export PILOT_LLM_MODEL=mistral:latest 
# 
# 
# PILOT_LLM_PROVIDER=ollama PILOT_LLM_MODEL=mistral:latest OLLAMA_URL=http://127.0.0.1:11434/api/chat PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak .venv/bin/python pilot_asr.py   


# export PILOT_MCP_COMMAND="node"
# export PILOT_MCP_ARGS="/Users/pagupta/Desktop/Capstone project1-2/index.js"
# export PILOT_MCP_TOOL="get_flights"

# # Start the pipeline
# PILOT_LLM_PROVIDER=ollama PILOT_LLM_MODEL=mistral:latest OLLAMA_URL=http://127.0.0.1:11434/api/chat PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak .venv/bin/python pilot_asr.py



# export PILOT_MCP_COMMAND="node"
# export PILOT_MCP_ARGS="/Users/pagupta/Desktop/Capstone_project1_2/flight-mcp/build/index.js"
# export PILOT_MCP_TOOL="get_flights"

# # Start the pipeline
# PILOT_LLM_PROVIDER=ollama PILOT_LLM_MODEL=mistral:latest OLLAMA_URL=http://127.0.0.1:11434/api/chat PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak .venv/bin/python pilot_asr.py



# export PILOT_MCP_COMMAND="node"
# export PILOT_MCP_ARGS="/Users/pagupta/Desktop/Capstone_project1_2/flight-mcp/index.js"
# export PILOT_MCP_TOOL="search_flights"

# PILOT_LLM_PROVIDER=ollama \
# PILOT_LLM_MODEL=mistral:latest \
# OLLAMA_URL=http://127.0.0.1:11434/api/chat \
# PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak \
# .venv/bin/python pilot_asr.py



# .venv/bin/python ppt_to_reveal.py ppt.pptx --output slides.html

# terminal 1
# .venv/bin/python slide_server.py 

# terminal 2

# export PILOT_MCP_COMMAND="node"
# export PILOT_MCP_ARGS="/Users/pagupta/Desktop/Capstone_project1_2/flight-mcp/index.js"
# export PILOT_MCP_TOOL="search_flights"
# PILOT_LLM_PROVIDER=ollama \
# PILOT_LLM_MODEL=mistral:latest \
# OLLAMA_URL=http://127.0.0.1:11434/api/chat \
# PILOT_TTS_SPEAK_URL=http://127.0.0.1:8000/speak \
# .venv/bin/python pilot_asr.py


# terminal 3
# .venv/bin/python -m uvicorn tts_server:app --reload
