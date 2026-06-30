"""
Unified PowerPoint API — upload (instant, stream, prepare), high-fidelity rendering,
navigation, Reveal.js viewer, and Slide Q&A.
Consolidates both ppt.py and ppt_instant.py to eliminate duplication.
"""


#                                       USER
#                                         │
#                      Upload PPT / Navigate / Ask Question
#                                         │
#                                         ▼
#                         ┌────────────────────────────────┐
#                         │      FastAPI PPT Router        │
#                         │        (Unified ppt.py)        │
#                         └────────────────────────────────┘
#                                         │
#           ┌─────────────────────────────┼─────────────────────────────┐
#           │                             │                             │
#           ▼                             ▼                             ▼
#    Upload APIs                  Viewer APIs                  AI & Navigation APIs
# ───────────────────      ───────────────────────      ───────────────────────────
# /upload                 /viewer/{session_id}         /navigate
# /upload_instant         /slides/{session_id}         /jump
# /upload_stream                                     /summarise
# /upload_prepare                                   /qa
# /render_stream
#           │
#           ▼
# ────────────────────────────────────────────────────────────────────────────
#                     PowerPoint Processing Pipeline
# ────────────────────────────────────────────────────────────────────────────
#           │
#           ▼
# Read PPT File
# (python-pptx)
#           │
#           ▼
# Presentation Object
#           │
#           ▼
# Loop Through Every Slide
#           │
#           ├─────────────────────────────────────┐
#           │                                     │
#           ▼                                     ▼
# Extract Metadata                      Render Slide Image
# (title, bullets, notes)                      │
#                                              │
#                             ┌────────────────┴────────────────┐
#                             │                                 │
#                             ▼                                 ▼
#                    LibreOffice Available?                  No
#                             │                                 │
#                      Yes ───┘                                 │
#                             ▼                                 ▼
#                 PPT → PDF Conversion             SVG Renderer (_slide_to_png_b64)
#                      (soffice)                             │
#                             │                             │
#                             ▼                             ▼
#                     PyMuPDF Reads PDF              PyMuPDF Renders SVG
#                             │                             │
#                             └──────────────┬──────────────┘
#                                            ▼
#                                Base64 PNG Image
#                                            │
#                                            ▼
#                              Slide Metadata + Image
#                                            │
#                                            ▼
#                      Store in _slide_store (Memory Cache)
#                                            │
#           ┌────────────────────────────────┼──────────────────────────────────┐
#           │                                │                                  │
#           ▼                                ▼                                  ▼
#  Reveal.js Viewer                  Navigation Engine                   AI Copilot
#           │                                │                                  │
#           ▼                                ▼                                  ▼
#  Browser HTML                  Event Bus / Commands              Summary / Q&A
# (next/prev/jump)                next, prev, goto                 using metadata

import asyncio
import base64
import html
import io
import json
import logging
import os
import re
import shutil
import subprocess
import tempfile
import time
from typing import List, Optional

from fastapi import APIRouter, File, HTTPException, UploadFile
from fastapi.responses import HTMLResponse, StreamingResponse
from pydantic import BaseModel

from backend.core.slide_store import (
    _current_slide,
    _slide_store,
    get_latest_upload_sid,
    resolve_sid,
    set_latest_upload_sid,
)

router = APIRouter()
logger = logging.getLogger("pilot.api.ppt")


# ============================================================
# MODELS
# ============================================================


class PPTCmd(BaseModel):
    session_id: str
    direction: str
    slide_index: int = -1


class JumpCmd(BaseModel):
    session_id: str
    query: str


class CreateCmd(BaseModel):
    session_id: str
    prompt: str
    slide_count: int = 5


# ============================================================
# RENDERERS & METADATA HELPERS
# ============================================================


def _slide_to_png_b64(slide, prs_w: int, prs_h: int, width_px: int = 960, height_px: int = 540) -> str:
    """Render one slide to a base64-encoded PNG string using PyMuPDF (no cairosvg/cairo dependency)."""
    import fitz

    scale_x = width_px / prs_w
    scale_y = height_px / prs_h

    # Background fill
    bg_color = "#111112"
    try:
        bg = slide.background.fill
        if hasattr(bg, "fore_color") and bg.fore_color.type:
            rgb = bg.fore_color.rgb
            bg_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
    except Exception:
        pass

    text_els = ""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        x = int(shape.left * scale_x)
        y = int(shape.top * scale_y)

        cursor_y = y + 22
        for para in shape.text_frame.paragraphs:
            raw = para.text.strip()
            if not raw:
                cursor_y += 12
                continue

            font_size, bold, color = 18, False, "#ffffff"
            try:
                if para.font.size:
                    font_size = int(para.font.size.pt)
                if para.font.bold:
                    bold = True
                if para.font.color and para.font.color.type:
                    rgb = para.font.color.rgb
                    color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
            except Exception:
                pass
            # A paragraph may contain mixed formatting
            for run in para.runs:
                try:
                    if run.font.size:
                        font_size = int(run.font.size.pt)
                    if run.font.bold:
                        bold = True
                    if run.font.color and run.font.color.type:
                        rgb = run.font.color.rgb
                        color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                except Exception:
                    pass

            fs = max(int(font_size * min(scale_x, scale_y) * 0.85), 8)
            fw = "bold" if bold else "normal"
            safe = html.escape(raw)
            # Wrap long text with tspan elements
            words = safe.split()
            lines, cur_line, max_chars = [], [], int(width_px / (fs * 0.6))
            for w in words:
                cur_line.append(w)
                if len(" ".join(cur_line)) > max_chars:
                    lines.append(" ".join(cur_line[:-1]))
                    cur_line = [w]
            if cur_line:
                lines.append(" ".join(cur_line))

            for line in lines:
                text_els += (
                    f'<text x="{x + 6}" y="{cursor_y}" '
                    f'font-size="{fs}" font-weight="{fw}" fill="{color}" '
                    f'font-family="Inter,Arial,sans-serif">'
                    f"{line}</text>\n"
                )
                cursor_y += int(fs * 1.45)
            cursor_y += 4

    svg = (
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}">'
        f'<rect width="{width_px}" height="{height_px}" fill="{bg_color}"/>'
        f"{text_els}"
        f"</svg>"
    )

    # Render SVG directly to PNG in memory using PyMuPDF (completely self-contained, no native cairo library needed)
    try:
        svg_doc = fitz.open("svg", svg.encode("utf-8"))
        page = svg_doc[0]
        pix = page.get_pixmap()
        png_data = pix.tobytes("png")
        svg_doc.close()
        return base64.b64encode(png_data).decode()
    except Exception as e:
        logger.error(f"Fallback SVG-to-PNG render failed: {e}")
        return ""


def _extract_slide_meta(slide, index: int, prs_w: int, prs_h: int) -> dict:
    """Fast, highly-structured metadata extraction (title, bullets, notes) for one slide.
    Filters out decorative elements, footer numbers, and topics to prevent duplicate compounding."""
    title = ""
    try:
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
    except Exception:
        pass

    # If title is still empty, look for a text shape at the top (top < 1.5 inches)
    if not title:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # pptx units are EMUs. 1 inch = 914400 EMUs.
                # 1.5 inches = 1371600 EMUs.
                if shape.top < 1371600:
                    text = shape.text_frame.text.strip()
                    if text and not any(k in text for k in ["Grid Dynamics", "PILOT Voice OS", "Topic:"]):
                        title = text
                        break

    bullets = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Skip if it is the title shape
            try:
                if slide.shapes.title == shape:
                    continue
            except Exception:
                pass

            # Skip footer/topic/header shapes
            shape_text = shape.text_frame.text.strip()
            if not shape_text:
                continue
            if shape_text == title:
                continue
            if (
                any(k in shape_text for k in ["Grid Dynamics", "PILOT Voice OS"])
                or shape_text.startswith("Topic:")
                or shape_text.startswith("Slide ")
            ):
                continue

            # Extract paragraphs as individual bullets
            for para in shape.text_frame.paragraphs:
                p_text = para.text.strip()
                if not p_text:
                    continue
                # Clean leading bullet symbols: •, *, -, etc.
                p_cleaned = re.sub(r"^[•\-\*\s·]+", "", p_text).strip()
                if p_cleaned and p_cleaned != title:
                    if (
                        not any(k in p_cleaned for k in ["Grid Dynamics", "PILOT Voice OS"])
                        and not p_cleaned.startswith("Topic:")
                        and not p_cleaned.startswith("Slide ")
                    ):
                        bullets.append(p_cleaned)

    notes = ""
    try:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()[:300]
    except Exception:
        pass

    return {
        "index": index,
        "title": title or f"Slide {index + 1}",
        "bullets": bullets[:10],
        "notes": notes,
    }


def _get_libreoffice_path() -> str | None:
    """Resolves headless soffice executable binary path across platforms."""
    soffice_path = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice_path:
        for p in [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/opt/homebrew/bin/soffice",
            "/usr/local/bin/soffice",
        ]:
            if os.path.exists(p):
                soffice_path = p
                break
    return soffice_path


def _build_reveal_html(slides: list, session_id: str) -> str:
    """Generates the full HTML markup with embedded slides for Reveal.js."""
    sections = ""
    for i, s in enumerate(slides):
        img_src = f"data:image/png;base64,{s['img_b64']}" if s.get("img_b64") else ""
        sections += f"""
        <section id="slide-{i}" style="text-align:center;">
            <img src="{img_src}" style="max-width:100%; max-height:100%; object-fit:contain; border:none; box-shadow:none; background:transparent;" />
        </section>
        """

    return f"""
<!DOCTYPE html>
<html>
<head>
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/reveal.min.css">
<style>
html, body {{
    margin:0;
    width:100%;
    height:100%;
    background:#111;
}}
.reveal {{
    width:100%;
    height:100%;
}}
</style>
</head>
<body>
<div class="reveal">
    <div class="slides">
        {sections}
    </div>
</div>
<script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/reveal.min.js"></script>
<script>
Reveal.initialize({{
    controls:true,
    progress:true,
    slideNumber:true,
    hash:false
}});

window.addEventListener("message", (e) => {{
    const action = e.data?.action;
    if(action === "next")
        Reveal.next();
    else if(action === "prev")
        Reveal.prev();
    else if(action === "first")
        Reveal.slide(0);
    else if(action === "last")
        Reveal.slide({len(slides) - 1});
    else if(action === "goto" && e.data.index !== undefined)
        Reveal.slide(e.data.index);
}});

Reveal.on("slidechanged", (event) => {{
    parent.postMessage({{
        type:"slide_changed",
        index:event.indexh
    }}, "*");
}});
</script>
</body>
</html>
"""


# ============================================================
# ENDPOINTS
# ============================================================

# ── Route A: Original upload (renders everything at once via PyMuPDF) ──


@router.post("/upload")
async def upload_ppt(session_id: str, file: UploadFile = File(...)):
    """Upload PPTX → returns simple status. Renders all slides synchronously via PyMuPDF."""
    if not file.filename.lower().endswith((".pptx", ".ppt")):
        raise HTTPException(status_code=400, detail="Only PPT/PPTX supported")

    content = await file.read()

    def _render():
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides):
            meta = _extract_slide_meta(slide, i, prs.slide_width, prs.slide_height)
            meta["img_b64"] = _slide_to_png_b64(slide, prs.slide_width, prs.slide_height)
            slides.append(meta)
        return slides

    slides = await asyncio.to_thread(_render)
    _slide_store[session_id] = slides
    set_latest_upload_sid(session_id)
    _current_slide[session_id] = 0

    return {"status": "ok", "slide_count": len(slides)}


# ── Route B: upload_instant (PDF-rendered slide flow with PyMuPDF fallback) ──


@router.post("/upload_instant")
async def upload_instant(session_id: str, file: UploadFile = File(...)):
    """Upload PPTX → returns JSON with all slides including high-fidelity images."""
    if not file.filename.lower().endswith((".pptx", ".ppt")):
        raise HTTPException(400, "Only .pptx / .ppt files supported")

    content = await file.read()

    def _render_all() -> list[dict]:
        import fitz
        from pptx import Presentation

        try:
            if hasattr(fitz, "TOOLS"):
                fitz.TOOLS.mupdf_display_errors(False)
        except Exception:
            pass

        prs = Presentation(io.BytesIO(content))
        prs_w = prs.slide_width
        prs_h = prs.slide_height

        slides_out = []
        for i, slide in enumerate(prs.slides):
            meta = _extract_slide_meta(slide, i, prs_w, prs_h)
            slides_out.append(meta)

        soffice_path = _get_libreoffice_path()
        pdf_rendered = False

        if soffice_path:
            try:
                with tempfile.TemporaryDirectory() as temp_dir:
                    pptx_temp_path = os.path.join(temp_dir, "presentation.pptx")
                    with open(pptx_temp_path, "wb") as f_temp:
                        f_temp.write(content)

                    cmd = [
                        soffice_path,
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        temp_dir,
                        pptx_temp_path,
                    ]
                    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

                    pdf_path = os.path.join(temp_dir, "presentation.pdf")
                    if os.path.exists(pdf_path):
                        doc = fitz.open(pdf_path)
                        for i, page in enumerate(doc):
                            if i < len(slides_out):
                                pix = page.get_pixmap(dpi=120)
                                png_bytes = pix.tobytes("png")
                                slides_out[i]["img_b64"] = base64.b64encode(png_bytes).decode()
                        doc.close()
                        pdf_rendered = True
                        logger.info("Successfully rendered all slides to PNG via LibreOffice PDF pipeline.")
            except Exception as lo_err:
                logger.error(f"LibreOffice instant conversion failed: {lo_err}")

        if not pdf_rendered:
            for i, slide in enumerate(prs.slides):
                slides_out[i]["img_b64"] = _slide_to_png_b64(slide, prs_w, prs_h)

        return slides_out

    slides = await asyncio.to_thread(_render_all)
    _slide_store[session_id] = slides

    return {"status": "ok", "slide_count": len(slides), "slides": slides}


# ── Route C: upload_stream (streams PNG slides slide-by-slide via SSE) ──


@router.post("/upload_stream")
async def upload_stream(session_id: str, file: UploadFile = File(...)):
    """Upload PPTX → SSE stream. Emits rendered PNG slides slide-by-slide."""
    if not file.filename.lower().endswith((".pptx", ".ppt")):
        raise HTTPException(400, "Only .pptx / .ppt files supported")

    content = await file.read()

    async def _event_generator():
        import fitz
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        prs_w = prs.slide_width
        prs_h = prs.slide_height
        total = len(prs.slides)
        store_list: list[dict] = []

        yield f"data: {json.dumps({'type': 'init', 'total': total})}\n\n"

        soffice_path = _get_libreoffice_path()
        pdf_doc = None

        if soffice_path:
            try:
                temp_dir = tempfile.mkdtemp()
                pptx_temp_path = os.path.join(temp_dir, "presentation.pptx")
                with open(pptx_temp_path, "wb") as f_temp:
                    f_temp.write(content)

                cmd = [
                    soffice_path,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    temp_dir,
                    pptx_temp_path,
                ]
                subprocess.run(
                    cmd, check=True, timeout=10.0, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
                )

                pdf_path = os.path.join(temp_dir, "presentation.pdf")
                if os.path.exists(pdf_path):
                    pdf_doc = fitz.open(pdf_path)
            except Exception as lo_err:
                logger.error(f"SSE background LibreOffice PDF conversion failed: {lo_err}")

        for i, slide in enumerate(prs.slides):
            meta = await asyncio.to_thread(_extract_slide_meta, slide, i, prs_w, prs_h)
            img_b64 = ""

            if pdf_doc and i < len(pdf_doc):
                try:
                    page = pdf_doc[i]
                    pix = await asyncio.to_thread(page.get_pixmap, dpi=120)
                    png_bytes = pix.tobytes("png")
                    img_b64 = base64.b64encode(png_bytes).decode()
                except Exception as page_err:
                    logger.error(f"Failed to render slide page {i} from PDF: {page_err}")

            if not img_b64:
                img_b64 = await asyncio.to_thread(_slide_to_png_b64, slide, prs_w, prs_h)

            store_list.append({**meta, "img_b64": img_b64})
            payload = {**meta, "img_b64": img_b64, "total": total, "type": "slide"}
            yield f"data: {json.dumps(payload)}\n\n"

        if pdf_doc:
            pdf_doc.close()
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except Exception:
                pass

        _slide_store[session_id] = store_list
        set_latest_upload_sid(session_id)
        _current_slide[session_id] = 0

        yield f"data: {json.dumps({'type': 'done', 'total': total})}\n\n"

    return StreamingResponse(
        _event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Access-Control-Allow-Origin": "*",
        },
    )


# ── Route D: Prepared stream (upload_prepare + render_stream split) ──


@router.post("/upload_prepare")
async def upload_prepare(session_id: str, file: UploadFile = File(...)):
    """Fast file saver to prepare the presentation file on disk."""
    if not file.filename.lower().endswith((".pptx", ".ppt")):
        raise HTTPException(400, "Only .pptx / .ppt files supported")

    content = await file.read()
    os.makedirs("data/ppt", exist_ok=True)
    path = f"data/ppt/{session_id}.pptx"
    with open(path, "wb") as f:
        f.write(content)
    return {"status": "ok"}


@router.get("/render_stream")
async def render_stream(session_id: str, token: str):
    """EventSource endpoint called by the frontend to render the slides page-by-page."""
    from backend.core.security import decode_token

    try:
        decode_token(token)
    except Exception:
        raise HTTPException(status_code=401, detail="Missing or invalid credentials")

    path = f"data/ppt/{session_id}.pptx"
    if not os.path.exists(path):
        raise HTTPException(404, "Prepared presentation file not found")

    async def _event_generator():
        import fitz
        from pptx import Presentation

        try:
            if hasattr(fitz, "TOOLS"):
                fitz.TOOLS.mupdf_display_errors(False)
        except Exception:
            pass

        prs = Presentation(path)
        prs_w = prs.slide_width
        prs_h = prs.slide_height
        total = len(prs.slides)
        store_list: list[dict] = []

        yield f"data: {json.dumps({'type': 'init', 'total': total})}\n\n"

        soffice_path = _get_libreoffice_path()
        pdf_doc = None

        if soffice_path:
            try:
                temp_dir = tempfile.mkdtemp()
                cmd = [soffice_path, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, path]
                subprocess.run(
                    cmd, check=True, timeout=15.0, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
                )

                pdf_filename = os.path.splitext(os.path.basename(path))[0] + ".pdf"
                pdf_path = os.path.join(temp_dir, pdf_filename)
                if os.path.exists(pdf_path):
                    pdf_doc = fitz.open(pdf_path)
            except Exception as lo_err:
                logger.error(f"render_stream LibreOffice PDF conversion failed: {lo_err}")

        for i, slide in enumerate(prs.slides):
            meta = await asyncio.to_thread(_extract_slide_meta, slide, i, prs_w, prs_h)
            img_b64 = ""

            if pdf_doc and i < len(pdf_doc):
                try:
                    page = pdf_doc[i]
                    pix = await asyncio.to_thread(page.get_pixmap, dpi=120)
                    png_bytes = pix.tobytes("png")
                    img_b64 = base64.b64encode(png_bytes).decode()
                except Exception as page_err:
                    logger.error(f"Failed to render slide page {i} from PDF: {page_err}")

            if not img_b64:
                img_b64 = await asyncio.to_thread(_slide_to_png_b64, slide, prs_w, prs_h)

            store_list.append({**meta, "img_b64": img_b64})
            payload = {**meta, "img_b64": img_b64, "total": total, "type": "slide"}
            yield f"data: {json.dumps(payload)}\n\n"

        if pdf_doc:
            pdf_doc.close()
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except Exception:
                pass

        _slide_store[session_id] = store_list
        set_latest_upload_sid(session_id)
        _current_slide[session_id] = 0

        yield f"data: {json.dumps({'type': 'done', 'total': total})}\n\n"

    return StreamingResponse(
        _event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Access-Control-Allow-Origin": "*",
        },
    )


# ── Route E: Viewer & Slide Control ──


@router.get("/viewer/{session_id}", response_class=HTMLResponse)
async def get_viewer(session_id: str):
    """Returns a full, independent Reveal.js HTML slide presentation page."""
    slides = _slide_store.get(session_id)
    if not slides:
        return HTMLResponse("<h2>No slides uploaded.</h2>")
    return HTMLResponse(_build_reveal_html(slides, session_id))


@router.post("/navigate")
async def navigate(cmd: PPTCmd):
    """Executes a slide navigation action (next, prev, first, last) on the active session."""
    from backend.tools.ppt_copilot import ppt_navigate

    return await ppt_navigate({"direction": cmd.direction}, cmd.session_id)


@router.post("/jump")
async def jump(cmd: JumpCmd):
    """Analyzes natural language queries and navigates to matching slide index numbers."""
    from backend.queues.bus import bus

    slides = _slide_store.get(cmd.session_id, [])
    q = cmd.query.lower()
    m = re.search(r"\b(\d+)\b", q)
    if m:
        idx = int(m.group(1)) - 1
        await bus.emit_event("ppt_command", {"action": "goto", "index": idx}, cmd.session_id)
        return {"status": "ok", "index": idx}
    return {"status": "ok", "index": 0}


# ── Route F: Slide Summary & Q&A ──


@router.post("/summarise")
async def summarise(data: dict):
    """Compiles the first 20 slide titles to provide a fast contextual synopsis of the slide deck."""
    slides = _slide_store.get(data.get("session_id", ""), [])
    titles = ", ".join(slide["title"] for slide in slides[:20])
    return {"reply": f"Presentation contains {len(slides)} slides. Topics: {titles}"}


@router.get("/slides/{session_id}")
async def get_slides(session_id: str):
    """Retrieves list of slide metadata and images for the frontend viewer state."""
    slides = _slide_store.get(session_id) or _slide_store.get(get_latest_upload_sid(), [])
    return {"slides": slides}


@router.post("/qa")
async def qa(data: dict):
    """Asks questions about the slide deck, invoking the local PDF/text shape parser."""
    from backend.tools.ppt_copilot import ppt_qa

    return await ppt_qa({"query": data.get("query", "")}, data.get("session_id", ""))


async def download_ai_image(prompt_str: str) -> str | None:
    """Helper to generate and download a professional AI image via Pollinations AI with rate-limit retries and fallbacks."""
    import asyncio
    import random
    import urllib.parse

    import httpx

    try:
        os.makedirs("data/ppt/images", exist_ok=True)
        filename = f"ai_gen_{int(time.time())}_{hash(prompt_str) % 10000}.png"
        save_path = os.path.join("data/ppt/images", filename)

        encoded_prompt = urllib.parse.quote(prompt_str)
        urls = [
            f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=800&height=600&model=flux&nologo=true",
            f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=800&height=600&nologo=true",
        ]

        async with httpx.AsyncClient(timeout=30.0) as client:
            for attempt in range(4):
                url = urls[0] if attempt < 2 else urls[1]
                try:
                    # Apply sleep with jitter to stagger parallel downloads
                    sleep_time = (attempt * 1.5) + random.uniform(0.5, 1.5)
                    await asyncio.sleep(sleep_time)

                    resp = await client.get(url)
                    if resp.status_code == 200:
                        with open(save_path, "wb") as f:
                            f.write(resp.content)
                        logger.info(f"Downloaded generated AI image to: {save_path} on attempt {attempt + 1}")
                        return save_path
                    elif resp.status_code == 429:
                        logger.warning(
                            f"Pollinations AI rate limited (429) for prompt: '{prompt_str}'. Retrying..."
                        )
                    else:
                        logger.warning(
                            f"Pollinations AI returned status {resp.status_code} for prompt: '{prompt_str}'"
                        )
                except Exception as attempt_err:
                    logger.error(f"Image download attempt {attempt + 1} failed: {attempt_err}")

        # Ultimate fallback: high-quality neutral professional placeholder image
        try:
            fallback_url = (
                "https://images.unsplash.com/photo-1557804506-669a67965ba0?auto=format&fit=crop&w=800&q=80"
            )
            async with httpx.AsyncClient(timeout=10.0) as client:
                resp = await client.get(fallback_url)
                if resp.status_code == 200:
                    with open(save_path, "wb") as f:
                        f.write(resp.content)
                    logger.info(f"Rate limited. Used high-quality fallback image for prompt: '{prompt_str}'")
                    return save_path
        except Exception as fb_err:
            logger.error(f"Ultimate fallback image download failed: {fb_err}")

    except Exception as e:
        logger.error(f"Failed to generate and download AI image: {e}")
    return None


def clear_ppt_directory():
    """Removes all files in data/ppt and data/ppt/images to start fresh."""
    import shutil

    try:
        ppt_dir = "data/ppt"
        if os.path.exists(ppt_dir):
            for filename in os.listdir(ppt_dir):
                file_path = os.path.join(ppt_dir, filename)
                try:
                    if os.path.isdir(file_path):
                        shutil.rmtree(file_path)
                    else:
                        os.unlink(file_path)
                except Exception as e:
                    logger.warning(f"Failed to delete {file_path}: {e}")
        os.makedirs(os.path.join(ppt_dir, "images"), exist_ok=True)
        logger.info("Successfully cleaned and initialized data/ppt directory.")
    except Exception as e:
        logger.error(f"Error clearing ppt directory: {e}")


async def create_presentation_from_prompt(prompt: str, session_id: str, slide_count: int = 5) -> list[dict]:
    """Generates slide outline using Ollama, creates a PPTX presentation, and renders the slides."""
    clear_ppt_directory()

    import httpx

    from backend.core.config import settings

    system_content = (
        "You are an expert presentation designer. Create a slide deck based on the user's topic.\n"
        "You MUST respond ONLY with a valid JSON object. Do not write any markdown, "
        "explanations, or backticks outside the JSON. The JSON must follow this exact schema:\n"
        "{\n"
        '  "slides": [\n'
        "    {\n"
        '      "title": "Slide Title",\n'
        '      "bullets": ["Key bullet point 1", "Key bullet point 2"],\n'
        '      "notes": "Detailed speaker notes for the presenter",\n'
        '      "image_prompt": "Detailed description of a professional graphic, chart, 3D object, or diagram that visually represents this slide\'s content, or \'None\' if an image is not suitable"\n'
        "    }\n"
        "  ]\n"
        "}\n"
        f"Generate exactly {slide_count} slides. Ensure the first slide is a high-impact Title Slide. "
        "Keep bullets concise and professional (under 12 words each, max 4 bullets per slide)."
    )

    url = f"{settings.OLLAMA_BASE_URL.rstrip('/')}/api/chat"
    payload = {
        "model": settings.OLLAMA_MODEL,
        "stream": False,
        "messages": [
            {"role": "system", "content": system_content},
            {"role": "user", "content": f"Create a presentation on the topic: {prompt}"},
        ],
        "options": {"temperature": 0.3, "num_predict": 1500},
        "format": "json",
    }

    slides_data = []
    parsed = None

    # 1. Try Gemini (Highly detailed, content-rich)
    if settings.GEMINI_API_KEY:
        try:
            logger.info("Generating content-rich slides using cloud Gemini API...")
            import google.generativeai as genai

            genai.configure(api_key=settings.GEMINI_API_KEY)
            model = genai.GenerativeModel("gemini-2.5-flash")

            prompt_content = (
                f"Create a highly detailed, professional presentation outline on the topic: '{prompt}'.\n"
                f"You MUST generate exactly {slide_count} slides.\n"
                "Ensure the slide content is extremely comprehensive, informative, and detailed (not just simple one-word bullet points).\n"
                "Respond ONLY with a valid JSON object matching this schema:\n"
                "{\n"
                '  "slides": [\n'
                "    {\n"
                '      "title": "Slide Title",\n'
                '      "bullets": ["Detailed, informative bullet point 1", "Detailed, informative bullet point 2"],\n'
                '      "notes": "Comprehensive speaker notes detailing the slide concepts",\n'
                '      "image_prompt": "Detailed description of a professional graphic, chart, 3D object, or diagram that visually represents this slide\'s content, or \'None\' if an image is not suitable"\n'
                "    }\n"
                "  ]\n"
                "}"
            )

            response = await model.generate_content_async(
                prompt_content, generation_config={"response_mime_type": "application/json"}
            )
            parsed = json.loads(response.text.strip())
        except Exception as gem_err:
            logger.warning(f"Cloud Gemini slide generation failed: {gem_err}. Falling back...")

    # 2. Try Groq (Llama 3.3 70B - ultra fast & detailed)
    if not parsed and settings.GROQ_API_KEY:
        try:
            logger.info("Generating content-rich slides using cloud Groq API...")
            from groq import AsyncGroq

            client = AsyncGroq(api_key=settings.GROQ_API_KEY)

            prompt_content = (
                f"Create a highly detailed, professional presentation outline on the topic: '{prompt}'.\n"
                f"You MUST generate exactly {slide_count} slides.\n"
                "Ensure the slide content is extremely comprehensive, informative, and detailed.\n"
                "Respond ONLY with a valid JSON object matching this schema:\n"
                "{\n"
                '  "slides": [\n'
                "    {\n"
                '      "title": "Slide Title",\n'
                '      "bullets": ["Detailed, informative bullet point 1", "Detailed, informative bullet point 2"],\n'
                '      "notes": "Comprehensive speaker notes detailing the slide concepts",\n'
                '      "image_prompt": "Detailed description of a professional graphic, chart, 3D object, or diagram that visually represents this slide\'s content, or \'None\' if an image is not suitable"\n'
                "    }\n"
                "  ]\n"
                "}"
            )

            resp = await client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert presentation designer. Respond only in JSON format matching the requested schema.",
                    },
                    {"role": "user", "content": prompt_content},
                ],
                response_format={"type": "json_object"},
            )
            parsed = json.loads(resp.choices[0].message.content)
        except Exception as groq_err:
            logger.warning(f"Cloud Groq slide generation failed: {groq_err}. Falling back...")

    # 3. Fallback to local Ollama
    if not parsed:
        try:
            logger.info(f"Generating slides using local Ollama model ({settings.OLLAMA_MODEL})...")
            async with httpx.AsyncClient(timeout=60.0) as client:
                resp = await client.post(url, json=payload)
                if resp.status_code == 200:
                    result_data = resp.json()
                    raw_content = result_data.get("message", {}).get("content", "").strip()

                    # Parse JSON
                    try:
                        parsed = json.loads(raw_content)
                    except Exception:
                        json_match = re.search(r"(\{.*\}|\[.*\])", raw_content, re.DOTALL)
                        if json_match:
                            parsed = json.loads(json_match.group(1))
                        else:
                            raise
        except Exception as ollama_err:
            logger.error(f"Local Ollama slide generation failed: {ollama_err}")

    # Normalize parsed JSON output into standard slide outline list
    try:
        if not parsed:
            raise ValueError("No slide data was successfully generated by any LLM provider.")

        # Robustly extract the list of slides
        extracted_list = []
        if isinstance(parsed, dict):
            # Find any list within the dictionary
            for key, val in parsed.items():
                if isinstance(val, list):
                    extracted_list = val
                    break
            if not extracted_list and "title" in parsed:
                extracted_list = [parsed]
        elif isinstance(parsed, list):
            extracted_list = parsed

        # Normalize each slide item in the list
        normalized_slides = []
        for i, item in enumerate(extracted_list):
            if isinstance(item, dict):
                bullets = item.get("bullets", [])
                if isinstance(bullets, str):
                    bullets = [bullets]
                elif not isinstance(bullets, list):
                    bullets = []
                bullets = [str(b) for b in bullets]

                # Retrieve the optional generated image prompt and download the AI asset
                image_prompt = item.get("image_prompt", "").strip()
                title_str = str(item.get("title", "Untitled Slide"))

                # Enforce visual asset generation for all content slides (slides after the title slide)
                if i > 0:
                    if not image_prompt or image_prompt.lower() in ("none", "null", ""):
                        image_prompt = (
                            f"A clean, professional modern corporate graphic representation of: {title_str}"
                        )

                slide_images = []
                if image_prompt and image_prompt.lower() not in ("none", "null", ""):
                    img_path = await download_ai_image(image_prompt)
                    if img_path:
                        slide_images.append(
                            {
                                "path": img_path,
                                "left_in": 7.5,
                                "top_in": 1.8,
                                "width_in": 5.0,
                                "height_in": 4.5,
                            }
                        )

                normalized_slides.append(
                    {
                        "title": title_str,
                        "bullets": bullets,
                        "notes": str(item.get("notes", "")),
                        "images": slide_images,
                    }
                )
            elif isinstance(item, str):
                normalized_slides.append({"title": item, "bullets": [], "notes": "", "images": []})

        if normalized_slides:
            slides_data = normalized_slides
        else:
            raise ValueError("No valid slides parsed from JSON")
    except Exception as e:
        logger.error(f"Failed to normalize slide data: {e}")
        # Fallback slides in case LLM generation fails
        slides_data = [
            {
                "title": f"Introduction to {prompt}",
                "bullets": [
                    f"Overview of {prompt}",
                    "Key concepts and foundations",
                    "Historical background and context",
                ],
                "notes": "Welcome everyone. Today we are exploring this topic.",
            },
            {
                "title": "Core Principles",
                "bullets": [
                    "Fundamental mechanism of action",
                    "Key structures and components",
                    "Important rules and guidelines",
                ],
                "notes": "Let's dive into the core mechanisms that make this work.",
            },
            {
                "title": "Main Benefits & Applications",
                "bullets": [
                    "Real-world use cases",
                    "Efficiency and cost improvements",
                    "Impact on modern industry",
                ],
                "notes": "Here are the primary ways this technology is applied in the real world.",
            },
            {
                "title": "Challenges & Solutions",
                "bullets": [
                    "Current limitations and hurdles",
                    "Innovative workarounds and research",
                    "Future outlook and developments",
                ],
                "notes": "Of course, there are some challenges we must address.",
            },
            {
                "title": "Conclusion & Q&A",
                "bullets": [
                    "Summary of key takeaways",
                    "Future roadmap and expansion",
                    "Thank you for your time",
                ],
                "notes": "Thank you. I am happy to open the floor for any questions.",
            },
        ]
        slides_data = slides_data[:slide_count]

    # Save and render presentation
    store_list = save_and_render_pptx(slides_data, session_id, prompt)
    _current_slide[session_id] = 0
    return store_list


def save_and_render_pptx(slides_data: list[dict], session_id: str, topic: str = "Presentation") -> list[dict]:
    """Compiles list of slide dictionaries to PPTX presentation and renders each slide to base64 PNG."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, slide_item in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # Widescreen left accent bar: left 0, top 0, width 0.12 inches, height 7.5 inches (Grid Dynamics Accent Bar)
        accent_bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE = 1
            Inches(0),
            Inches(0),
            Inches(0.12),
            Inches(7.5),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(245, 167, 0)  # Amber Gold (#F5A700)
        accent_bar.line.fill.background()

        # Background solid color: White (#FFFFFF)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        title_text = slide_item.get("title", f"Slide {idx + 1}")
        bullets = slide_item.get("bullets", [])
        notes = slide_item.get("notes", "")

        # Check if this slide has any valid images on disk
        has_images = False
        images = slide_item.get("images", [])
        if images and isinstance(images, list):
            for img_item in images:
                img_path = img_item.get("path")
                if img_path and os.path.exists(img_path):
                    has_images = True
                    break

        if idx == 0:
            # Title slide layout
            tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.0))
            tf = tx_box.text_frame
            tf.word_wrap = True

            # Title text
            p = tf.paragraphs[0]
            p.text = title_text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(17, 17, 18)  # Dark Charcoal

            # Subtitle (Darker Amber Gold for contrast)
            p2 = tf.add_paragraph()
            p2.text = bullets[0] if bullets else "Grid Dynamics Presentation Workspace"
            p2.alignment = PP_ALIGN.CENTER
            p2.font.name = "Arial"
            p2.font.size = Pt(22)
            p2.font.color.rgb = RGBColor(212, 144, 15)  # Dark Amber Gold
            p2.space_before = Pt(20)

            # Extra details (Medium Gray)
            if bullets and len(bullets) > 1:
                p3 = tf.add_paragraph()
                p3.text = " · ".join(bullets[1:])
                p3.alignment = PP_ALIGN.CENTER
                p3.font.name = "Arial"
                p3.font.size = Pt(14)
                p3.font.color.rgb = RGBColor(110, 110, 115)  # Medium Gray
                p3.space_before = Pt(14)
        else:
            # Content slide layout
            # Title at the top-left (Darker Amber Gold for contrast)
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(1.0))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            p_title.font.name = "Arial"
            p_title.font.size = Pt(32)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(212, 144, 15)  # Dark Amber Gold

            # Divider line decorative element (Amber Gold)
            divider = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE = 1
                Inches(0.8),
                Inches(1.4),
                Inches(11.733),
                Inches(0.02),
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = RGBColor(245, 167, 0)  # Amber Gold (#F5A700)
            divider.line.fill.background()

            # Bullet points — narrow width if slide contains an image
            content_width = Inches(6.2) if has_images else Inches(11.733)
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), content_width, Inches(4.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            # Limit bullets to at most 5 items to guarantee clean layout boundaries
            visible_bullets = bullets[:5]
            num_bullets = len(visible_bullets)
            total_chars = sum(len(b) for b in visible_bullets)

            # Dynamically compute optimal font size and spacing based on copy volume
            if num_bullets >= 5 or total_chars > 320:
                font_size = Pt(13)
                space_after = Pt(6)
            elif num_bullets == 4 or total_chars > 220:
                font_size = Pt(15)
                space_after = Pt(8)
            elif num_bullets == 3:
                font_size = Pt(17)
                space_after = Pt(10)
            else:
                font_size = Pt(20)
                space_after = Pt(14)

            for b_idx, b in enumerate(visible_bullets):
                p_bullet = tf_content.paragraphs[0] if b_idx == 0 else tf_content.add_paragraph()
                p_bullet.text = f"•  {b}"
                p_bullet.font.name = "Arial"
                p_bullet.font.size = font_size
                p_bullet.font.color.rgb = RGBColor(40, 40, 42)  # Dark Charcoal
                p_bullet.space_after = space_after

            # Insert images on the right side if present
            if has_images:
                for img_item in images:
                    img_path = img_item.get("path")
                    if img_path and os.path.exists(img_path):
                        try:
                            left_val = img_item.get("left_in", 7.5)
                            top_val = img_item.get("top_in", 1.8)
                            w_val = img_item.get("width_in", 5.0)
                            h_val = img_item.get("height_in", 4.5)
                            slide.shapes.add_picture(
                                img_path,
                                Inches(left_val),
                                Inches(top_val),
                                width=Inches(w_val),
                                height=Inches(h_val),
                            )
                        except Exception as img_err:
                            logger.error(f"Failed to add picture {img_path} to slide {idx + 1}: {img_err}")

            # Subtle footer divider
            footer_line = slide.shapes.add_shape(1, Inches(0.8), Inches(6.7), Inches(11.733), Inches(0.015))
            footer_line.fill.solid()
            footer_line.fill.fore_color.rgb = RGBColor(220, 220, 225)  # Light Gray
            footer_line.line.fill.background()

            # Slide number + branded footer (Grid Dynamics)
            footer_box = slide.shapes.add_textbox(Inches(8.0), Inches(6.8), Inches(4.533), Inches(0.5))
            p_foot = footer_box.text_frame.paragraphs[0]
            p_foot.text = f"Grid Dynamics  |  Slide {idx + 1} of {len(slides_data)}"
            p_foot.alignment = PP_ALIGN.RIGHT
            p_foot.font.name = "Arial"
            p_foot.font.size = Pt(9)
            p_foot.font.color.rgb = RGBColor(110, 110, 115)  # Medium Gray

            # Topic on bottom left
            topic_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(6.0), Inches(0.5))
            p_topic = topic_box.text_frame.paragraphs[0]
            p_topic.text = f"Topic: {topic.title()}"
            p_topic.font.name = "Arial"
            p_topic.font.size = Pt(9)
            p_topic.font.color.rgb = RGBColor(142, 142, 147)  # Warm Gray

        # Set speaker notes
        if notes:
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

    # Save presentation
    os.makedirs("data/ppt", exist_ok=True)
    pptx_path = f"data/ppt/{session_id}.pptx"
    prs.save(pptx_path)

    # Render the presentation to PNGs
    import fitz

    try:
        if hasattr(fitz, "TOOLS"):
            fitz.TOOLS.mupdf_display_errors(False)
    except Exception:
        pass

    prs_w = prs.slide_width
    prs_h = prs.slide_height
    store_list = []

    soffice_path = _get_libreoffice_path()
    pdf_doc = None

    if soffice_path:
        try:
            temp_dir = tempfile.mkdtemp()
            cmd = [soffice_path, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, pptx_path]
            subprocess.run(
                cmd, check=True, timeout=15.0, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
            )

            pdf_filename = f"{session_id}.pdf"
            pdf_path = os.path.join(temp_dir, pdf_filename)
            if os.path.exists(pdf_path):
                pdf_doc = fitz.open(pdf_path)
        except Exception as lo_err:
            logger.error(f"AI presentation LibreOffice PDF conversion failed: {lo_err}")

    for i, slide in enumerate(prs.slides):
        meta = _extract_slide_meta(slide, i, prs_w, prs_h)
        img_b64 = ""

        if pdf_doc and i < len(pdf_doc):
            try:
                page = pdf_doc[i]
                pix = page.get_pixmap(dpi=120)
                png_bytes = pix.tobytes("png")
                img_b64 = base64.b64encode(png_bytes).decode()
            except Exception as page_err:
                logger.error(f"Failed to render generated slide page {i} from PDF: {page_err}")

        if not img_b64:
            img_b64 = _slide_to_png_b64(slide, prs_w, prs_h)

        # Preserve the images data structure in memory store
        slide_images = slides_data[i].get("images", []) if i < len(slides_data) else []
        store_list.append({**meta, "img_b64": img_b64, "images": slide_images})

    if pdf_doc:
        pdf_doc.close()
        try:
            shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

    _slide_store[session_id] = store_list
    set_latest_upload_sid(session_id)
    return store_list


async def improvise_slide_content(current_slide: dict, user_prompt: str) -> dict:
    """Invokes LLM (Gemini -> Groq -> local Ollama cascade) to rewrite the content of a single slide."""
    import httpx

    from backend.core.config import settings

    system_content = (
        "You are an expert presentation improver. Your job is to rewrite the content of a SINGLE slide based on the user's improvisation request.\n"
        "Ensure the slide content is extremely comprehensive, informative, highly detailed, and professional (avoid simple short one-word or brief bullet points). "
        "Each bullet point must be a well-structured, complete, multi-sentence statement (at least 20-30 words) containing specific facts, concrete data, or detailed reasoning.\n"
        "IMPORTANT: Even if the user's request is to 'simplify', 'summarize', 'make concise', or 'shorten', you MUST NOT produce brief, low-content, or vague bullet points. "
        "Instead, rewrite using clear and highly readable language, but maintain a high density of information and a good amount of content—each bullet point must remain a fully developed, detailed sentence or two explaining a complete concept. Never output short phrases or single-word bullets.\n"
        "You MUST respond ONLY with a valid JSON object matching this exact schema:\n"
        "{\n"
        '  "title": "Improved Slide Title",\n'
        '  "bullets": ["Highly detailed, informative, and comprehensive bullet point 1", "Highly detailed, informative, and comprehensive bullet point 2"],\n'
        '  "notes": "Comprehensive and detailed speaker notes fully detailing the concepts shown on the slide"\n'
        "}\n"
        "Do not write any markdown, explanations, or backticks outside the JSON.\n"
        "Generate 3-5 comprehensive and rich bullet points that fully cover the subject matter."
    )

    current_desc = (
        f"Current Slide Content:\n"
        f"Title: {current_slide.get('title', 'Untitled')}\n"
        f"Bullets: {current_slide.get('bullets', [])}\n"
        f"Notes: {current_slide.get('notes', '')}\n"
    )

    user_content = (
        f"{current_desc}\n"
        f"Improvisation Request: {user_prompt}\n\n"
        f"Please rewrite the slide according to the request."
    )

    parsed = None

    # 1. Try Gemini
    if settings.GEMINI_API_KEY:
        try:
            import google.generativeai as genai

            genai.configure(api_key=settings.GEMINI_API_KEY)
            model = genai.GenerativeModel("gemini-2.5-flash")
            response = await model.generate_content_async(
                user_content,
                generation_config={
                    "response_mime_type": "application/json",
                    "system_instruction": system_content,
                },
            )
            parsed = json.loads(response.text.strip())
        except Exception as e:
            logger.warning(f"Cloud Gemini slide improvisation failed: {e}")

    # 2. Try Groq
    if not parsed and settings.GROQ_API_KEY:
        try:
            from groq import AsyncGroq

            client = AsyncGroq(api_key=settings.GROQ_API_KEY)
            resp = await client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[
                    {"role": "system", "content": system_content},
                    {"role": "user", "content": user_content},
                ],
                response_format={"type": "json_object"},
            )
            parsed = json.loads(resp.choices[0].message.content)
        except Exception as e:
            logger.warning(f"Cloud Groq slide improvisation failed: {e}")

    # 3. Fallback to local Ollama
    if not parsed:
        try:
            url = f"{settings.OLLAMA_BASE_URL.rstrip('/')}/api/chat"
            payload = {
                "model": settings.OLLAMA_MODEL,
                "stream": False,
                "messages": [
                    {"role": "system", "content": system_content},
                    {"role": "user", "content": user_content},
                ],
                "options": {"temperature": 0.3, "num_predict": 1500},
                "format": "json",
            }
            async with httpx.AsyncClient(timeout=40.0) as client:
                resp = await client.post(url, json=payload)
                if resp.status_code == 200:
                    result_data = resp.json()
                    raw_content = result_data.get("message", {}).get("content", "").strip()
                    try:
                        parsed = json.loads(raw_content)
                    except Exception:
                        json_match = re.search(r"(\{.*\})", raw_content, re.DOTALL)
                        if json_match:
                            parsed = json.loads(json_match.group(1))
        except Exception as e:
            logger.error(f"Local Ollama slide improvisation failed: {e}")

    # Normalize output
    if parsed and isinstance(parsed, dict):
        bullets = parsed.get("bullets", [])
        if isinstance(bullets, str):
            bullets = [bullets]
        elif not isinstance(bullets, list):
            bullets = []
        return {
            "title": str(parsed.get("title", current_slide.get("title", "Untitled"))),
            "bullets": [str(b) for b in bullets],
            "notes": str(parsed.get("notes", current_slide.get("notes", ""))),
        }

    # Unmodified fallback
    return {
        "title": current_slide.get("title", "Untitled"),
        "bullets": current_slide.get("bullets", []),
        "notes": current_slide.get("notes", ""),
    }


class AddSlideCmd(BaseModel):
    session_id: str
    title: Optional[str] = None
    bullets: Optional[List[str]] = None
    notes: Optional[str] = None


class EditSlideCmd(BaseModel):
    session_id: str
    slide_index: int
    title: Optional[str] = None
    bullets: Optional[List[str]] = None
    add_bullet: Optional[str] = None
    notes: Optional[str] = None
    images: Optional[List[dict]] = None


class ClearCmd(BaseModel):
    session_id: str


class DeleteSlideCmd(BaseModel):
    session_id: str
    slide_index: int


class ImproviseSlideCmd(BaseModel):
    session_id: str
    slide_index: int
    prompt: str


@router.post("/add_slide")
async def add_slide(cmd: AddSlideCmd):
    """Appends a slide manually or programmatically, saves & renders, and emits reload."""
    from backend.queues.bus import bus

    session_id = resolve_sid(cmd.session_id)
    slides = list(_slide_store.get(session_id, []))

    new_slide = {
        "title": cmd.title or "New Slide",
        "bullets": cmd.bullets or ["Write content or use voice to edit this slide"],
        "notes": cmd.notes or "",
        "images": [],
    }

    # Re-build plain slide data preserving images
    slides_to_save = []
    for s in slides:
        slides_to_save.append(
            {
                "title": s.get("title", ""),
                "bullets": s.get("bullets", []),
                "notes": s.get("notes", ""),
                "images": s.get("images", []),
            }
        )
    slides_to_save.append(new_slide)

    try:
        rendered = save_and_render_pptx(slides_to_save, session_id, topic="Presentation")
        new_idx = len(rendered) - 1
        _current_slide[session_id] = new_idx

        # Emit a reload event and navigate to the new slide
        await bus.emit_event(
            "ppt_command",
            {
                "action": "reload",
                "slides": rendered,
                "filename": f"Presentation ({len(rendered)} slides)",
                "index": new_idx,
                "preserveCurrent": False,
            },
            session_id,
        )

        return {"status": "ok", "slide_count": len(rendered), "slides": rendered, "index": new_idx}
    except Exception as e:
        logger.error(f"Error adding slide: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/delete_slide")
async def delete_slide(cmd: DeleteSlideCmd):
    """Deletes a slide by index, saves & renders, and emits reload."""
    from backend.queues.bus import bus

    session_id = resolve_sid(cmd.session_id)
    slides = list(_slide_store.get(session_id, []))
    if not slides:
        raise HTTPException(status_code=400, detail="No active presentation to delete slide from.")

    if cmd.slide_index < 0 or cmd.slide_index >= len(slides):
        raise HTTPException(
            status_code=400, detail=f"Slide index {cmd.slide_index} out of range (0-{len(slides) - 1})"
        )

    # Remove slide
    slides.pop(cmd.slide_index)

    # Re-build slide list preserving images
    slides_to_save = []
    for s in slides:
        slides_to_save.append(
            {
                "title": s.get("title", ""),
                "bullets": s.get("bullets", []),
                "notes": s.get("notes", ""),
                "images": s.get("images", []),
            }
        )

    try:
        rendered = save_and_render_pptx(slides_to_save, session_id, topic="Presentation")
        new_idx = min(cmd.slide_index, max(len(rendered) - 1, 0))
        _current_slide[session_id] = new_idx

        # Emit a reload event to update the frontend state in real-time
        await bus.emit_event(
            "ppt_command",
            {
                "action": "reload",
                "slides": rendered,
                "filename": f"Presentation ({len(rendered)} slides)",
                "index": new_idx,
                "preserveCurrent": False,
            },
            session_id,
        )

        return {"status": "ok", "slide_count": len(rendered), "slides": rendered, "index": new_idx}
    except Exception as e:
        logger.error(f"Error deleting slide: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/edit_slide")
async def edit_slide(cmd: EditSlideCmd):
    """Edits a slide by index, saves & renders, and emits reload preserving current index."""
    from backend.queues.bus import bus

    session_id = resolve_sid(cmd.session_id)
    slides = list(_slide_store.get(session_id, []))
    if not slides:
        raise HTTPException(status_code=400, detail="No active presentation to edit.")

    idx = cmd.slide_index
    if idx < 0 or idx >= len(slides):
        raise HTTPException(status_code=400, detail=f"Slide index {idx} out of range (0-{len(slides) - 1})")

    # Apply modifications
    if cmd.title is not None:
        slides[idx]["title"] = cmd.title
    if cmd.bullets is not None:
        slides[idx]["bullets"] = cmd.bullets
    elif cmd.add_bullet is not None:
        if "bullets" not in slides[idx] or not isinstance(slides[idx]["bullets"], list):
            slides[idx]["bullets"] = []
        slides[idx]["bullets"].append(cmd.add_bullet)
    if cmd.notes is not None:
        slides[idx]["notes"] = cmd.notes
    if cmd.images is not None:
        slides[idx]["images"] = cmd.images

    # Re-build plain slide data preserving images
    slides_to_save = []
    for s in slides:
        slides_to_save.append(
            {
                "title": s.get("title", ""),
                "bullets": s.get("bullets", []),
                "notes": s.get("notes", ""),
                "images": s.get("images", []),
            }
        )

    try:
        rendered = save_and_render_pptx(slides_to_save, session_id, topic="Presentation")
        _current_slide[session_id] = idx

        # Emit a reload event preserving current slide view
        await bus.emit_event(
            "ppt_command",
            {
                "action": "reload",
                "slides": rendered,
                "filename": f"Presentation ({len(rendered)} slides)",
                "index": idx,
                "preserveCurrent": True,
            },
            session_id,
        )

        return {"status": "ok", "slide_count": len(rendered), "slides": rendered, "index": idx}
    except Exception as e:
        logger.error(f"Error editing slide: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/clear")
async def clear_presentation_endpoint(cmd: ClearCmd):
    """Clears the presentation and starts with 1 blank title slide."""
    from backend.queues.bus import bus

    session_id = cmd.session_id
    default_slides = [
        {
            "title": "New Presentation",
            "bullets": ["Voice-Driven Interactive Presentation", "Start speaking or clicking to add slides"],
            "notes": "Welcome to your new interactive presentation deck.",
            "images": [],
        }
    ]

    try:
        rendered = save_and_render_pptx(default_slides, session_id, topic="New Presentation")
        _current_slide[session_id] = 0

        await bus.emit_event(
            "ppt_command",
            {
                "action": "reload",
                "slides": rendered,
                "filename": "New Presentation",
                "index": 0,
                "preserveCurrent": False,
            },
            session_id,
        )

        return {"status": "ok", "slide_count": len(rendered), "slides": rendered, "index": 0}
    except Exception as e:
        logger.error(f"Error clearing presentation: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/upload_image")
async def upload_image(session_id: str, slide_index: int, file: UploadFile = File(...)):
    """Uploads an image for a specific slide, saves it, inserts it into the slide design, and re-renders."""
    from backend.queues.bus import bus

    session_id = resolve_sid(session_id)
    slides = list(_slide_store.get(session_id, []))
    if not slides:
        raise HTTPException(status_code=400, detail="No active presentation to add image to.")

    if slide_index < 0 or slide_index >= len(slides):
        raise HTTPException(
            status_code=400, detail=f"Slide index {slide_index} out of range (0-{len(slides) - 1})"
        )

    # Ensure images directory exists
    img_dir = "data/ppt/images"
    os.makedirs(img_dir, exist_ok=True)

    # Save the file with a clean unique name
    ext = os.path.splitext(file.filename)[1] or ".png"
    filename = f"{session_id}_{slide_index}_{int(time.time())}{ext}"
    saved_path = os.path.join(img_dir, filename)

    try:
        content = await file.read()
        with open(saved_path, "wb") as f:
            f.write(content)

        # Add to slide metadata (clear previous images for a clean layout)
        slides[slide_index]["images"] = [
            {"path": saved_path, "left_in": 7.5, "top_in": 1.8, "width_in": 5.0, "height_in": 4.5}
        ]

        # Re-build plain slide data preserving images
        slides_to_save = []
        for s in slides:
            slides_to_save.append(
                {
                    "title": s.get("title", ""),
                    "bullets": s.get("bullets", []),
                    "notes": s.get("notes", ""),
                    "images": s.get("images", []),
                }
            )

        rendered = save_and_render_pptx(slides_to_save, session_id, topic="Presentation")
        _current_slide[session_id] = slide_index

        # Emit a reload event preserving current slide view
        await bus.emit_event(
            "ppt_command",
            {
                "action": "reload",
                "slides": rendered,
                "filename": f"Presentation ({len(rendered)} slides)",
                "index": slide_index,
                "preserveCurrent": True,
            },
            session_id,
        )

        return {"status": "ok", "slide_count": len(rendered), "slides": rendered, "index": slide_index}
    except Exception as e:
        logger.error(f"Error uploading image to slide: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/improvise_slide")
async def improvise_slide(cmd: ImproviseSlideCmd):
    """Uses LLM to improvise the content of a single slide based on user instructions, then saves and renders."""
    from backend.queues.bus import bus

    session_id = resolve_sid(cmd.session_id)
    slides = list(_slide_store.get(session_id, []))
    if not slides:
        raise HTTPException(status_code=400, detail="No active presentation to improvise.")

    idx = cmd.slide_index
    if idx < 0 or idx >= len(slides):
        raise HTTPException(status_code=400, detail=f"Slide index {idx} out of range (0-{len(slides) - 1})")

    current_slide = slides[idx]

    try:
        # Call LLM to improve slide
        improved = await improvise_slide_content(current_slide, cmd.prompt)

        # Update text fields but preserve any images!
        slides[idx]["title"] = improved["title"]
        slides[idx]["bullets"] = improved["bullets"]
        slides[idx]["notes"] = improved["notes"]

        # Re-build plain slide data preserving images
        slides_to_save = []
        for s in slides:
            slides_to_save.append(
                {
                    "title": s.get("title", ""),
                    "bullets": s.get("bullets", []),
                    "notes": s.get("notes", ""),
                    "images": s.get("images", []),
                }
            )

        rendered = save_and_render_pptx(slides_to_save, session_id, topic="Presentation")
        _current_slide[session_id] = idx

        # Emit a reload event preserving current slide view
        await bus.emit_event(
            "ppt_command",
            {
                "action": "reload",
                "slides": rendered,
                "filename": f"Presentation ({len(rendered)} slides)",
                "index": idx,
                "preserveCurrent": True,
            },
            session_id,
        )

        return {"status": "ok", "slide_count": len(rendered), "slides": rendered, "index": idx}
    except Exception as e:
        logger.error(f"Error improvising slide: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/create")
async def create_presentation(cmd: CreateCmd):
    """Generates a presentation dynamically from a text prompt using Ollama qwen2.5:7b, then renders it."""
    from backend.queues.bus import bus

    try:
        slides = await create_presentation_from_prompt(
            prompt=cmd.prompt, session_id=cmd.session_id, slide_count=cmd.slide_count
        )
        # Emit a reload event to update the frontend state in real-time
        await bus.emit_event(
            "ppt_command",
            {"action": "reload", "slides": slides, "filename": f"AI: {cmd.prompt}"},
            cmd.session_id,
        )

        return {"status": "ok", "slide_count": len(slides), "slides": slides}
    except Exception as e:
        logger.error(f"Error generating presentation: {e}")
        raise HTTPException(status_code=500, detail=str(e))
